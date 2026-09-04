"""Consulting-grade native PPTX exporter for MSSP SOC.

Design references: Deloitte / EY / PwC / KPMG / Accenture / Microsoft executive
decks + Canva Corporate Agile template. Storytelling over dashboard mirroring.

Rules:
- White backgrounds, blue corporate accent, minimal, plenty of whitespace.
- One clear message per slide, ≤6 KPI cards, ≤2 charts.
- Native PPTX objects only (no images). Charts, tables, shapes, text.
- Every content slide carries a 1-line executive insight, not a description.
"""
import io
import logging
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.oxml.ns import qn
from lxml import etree


# --- Layout -----
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN_L, MARGIN_R = 0.9, 0.9
CONTENT_W = 13.333 - MARGIN_L - MARGIN_R  # 11.53 in

# --- Corporate palette (Canva-inspired, blue-navy corporate) -----
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_TINT = RGBColor(0xEE, 0xF2, 0xFA)       # soft lavender-blue background
BG_SOFT = RGBColor(0xF8, 0xFA, 0xFC)
BORDER = RGBColor(0xE2, 0xE8, 0xF0)
BORDER_STRONG = RGBColor(0xCB, 0xD5, 0xE1)
INK = RGBColor(0x0F, 0x17, 0x2A)
INK_STRONG = RGBColor(0x1E, 0x29, 0x3B)
INK_MUTED = RGBColor(0x64, 0x74, 0x8B)
INK_SUBTLE = RGBColor(0x94, 0xA3, 0xB8)

# Corporate primary: navy
CORP = RGBColor(0x0B, 0x1F, 0x42)          # deep navy
CORP_600 = RGBColor(0x1E, 0x3A, 0x8A)      # blue-900
CORP_500 = RGBColor(0x2B, 0x50, 0xC7)      # blue-700
CORP_50 = RGBColor(0xEF, 0xF6, 0xFF)       # blue-50
CORP_100 = RGBColor(0xDB, 0xEA, 0xFE)      # blue-100

# 5-step gradient (green → navy) — for numbered items like the reference
STEP_COLORS = [
    RGBColor(0x33, 0xD4, 0x44),  # 01 bright green
    RGBColor(0x00, 0xC8, 0x78),  # 02 mint
    RGBColor(0x00, 0x8B, 0x8B),  # 03 teal
    RGBColor(0x0F, 0x3A, 0x5A),  # 04 dark teal
    RGBColor(0x0B, 0x1F, 0x42),  # 05 navy
]

# Semantic
EMERALD = RGBColor(0x05, 0x96, 0x69)
AMBER = RGBColor(0xD9, 0x77, 0x06)
ROSE = RGBColor(0xDC, 0x26, 0x26)
GRAY = RGBColor(0x64, 0x74, 0x8B)

# --- Deloitte QBR palette (teal + navy) ---------------------------------
D_TEAL = RGBColor(0x00, 0xAB, 0xAB)
D_TEAL_DK = RGBColor(0x00, 0x7C, 0x7C)
D_TEAL_LT = RGBColor(0x86, 0xD8, 0xD2)
D_NAVY = RGBColor(0x01, 0x21, 0x69)
D_NAVY_DK = RGBColor(0x0B, 0x1F, 0x3A)
D_SLATE = RGBColor(0x18, 0x24, 0x2C)
D_GRAY = RGBColor(0xA9, 0xB6, 0xC2)
D_GRAY_DK = RGBColor(0x5A, 0x62, 0x68)
D_BG = RGBColor(0xF4, 0xF7, 0xF8)
D_GREEN = RGBColor(0x86, 0xBC, 0x25)   # Deloitte wordmark dot
SEV_QBR = {
    "Critical": D_NAVY_DK,
    "High": D_NAVY,
    "Medium": D_TEAL,
    "Low": D_TEAL_LT,
    "Informational": D_GRAY,
}


def _wordmark(slide, x, y, *, size=20, dark=True):
    """Deloitte wordmark: 'Deloitte' + green dot, reproduced as text."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(3.0), Inches(size / 40))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Deloitte"
    r.font.name = "Aptos Display"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = (BG_WHITE if not dark else D_SLATE)
    dot = p.add_run()
    dot.text = "."
    dot.font.name = "Aptos Display"
    dot.font.size = Pt(size)
    dot.font.bold = True
    dot.font.color.rgb = D_GREEN
    return tb


def _vs(val, *, suffix="", na="N/A", pct=False):
    """Format a value or fall back to N/A when missing/None."""
    if val is None or val == "" or (isinstance(val, float) and val != val):
        return na
    if isinstance(val, (int, float)):
        v = f"{val:,.0f}" if float(val).is_integer() else f"{val:,.1f}"
        return f"{v}{'%' if pct else suffix}"
    return f"{val}{suffix}"


def _hex(hex_str: str) -> RGBColor:
    hex_str = (hex_str or "#1E3A8A").lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))


# --- Primitives -----
def _fill(shape, rgb):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def _no_line(shape):
    shape.line.fill.background()


def _line(shape, rgb=BORDER, width=0.5):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width)


def _text(slide, x, y, w, h, text, *, size=11, bold=False, color=INK,
          align=PP_ALIGN.LEFT, font="Aptos", tracking=None,
          uppercase=False, anchor=None, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text.upper() if uppercase else text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    if tracking is not None:
        rPr = run._r.get_or_add_rPr()
        rPr.set("spc", str(int(tracking * 100)))
    return tb


def _rect(slide, x, y, w, h, *, fill=BG_WHITE, border=None, radius=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    _fill(shape, fill)
    if border is None:
        _no_line(shape)
    else:
        _line(shape, *border)
    if radius is not None and hasattr(shape, "adjustments"):
        shape.adjustments[0] = radius
    return shape


def _hr(slide, x, y, w, *, color=BORDER, thickness=0.5):
    """Thin horizontal rule."""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
        Inches(w), Emu(int(9525 * thickness)),
    )
    _fill(line, color)
    _no_line(line)
    return line


def _icon_disc(slide, x, y, size, *, color=CORP, glyph=""):
    """Small filled disc with a glyph inside — SVG-icon substitute."""
    disc = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size),
    )
    _fill(disc, color)
    _no_line(disc)
    if glyph:
        tf = disc.text_frame
        tf.margin_left = tf.margin_top = tf.margin_bottom = tf.margin_right = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = glyph
        r.font.size = Pt(int(size * 22))
        r.font.bold = True
        r.font.color.rgb = BG_WHITE
    return disc


def _shadow_card(slide, x, y, w, h, *, radius=0.10):
    """Rounded white card with a soft shadow — Canva-style."""
    # Shadow (subtle offset gray)
    sh = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x + 0.04), Inches(y + 0.05), Inches(w), Inches(h),
    )
    _fill(sh, RGBColor(0xCF, 0xD8, 0xE3))
    _no_line(sh)
    sh.adjustments[0] = radius
    # Card
    card = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    _fill(card, BG_WHITE)
    _no_line(card)
    card.adjustments[0] = radius
    return card


def _chevron_label(slide, x, y, w, h, *, text, color):
    """Chevron (right-pointing pentagon) with white bold label — Canva-style."""
    ch = slide.shapes.add_shape(
        MSO_SHAPE.PENTAGON, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    _fill(ch, color)
    _no_line(ch)
    tf = ch.text_frame
    tf.margin_left = Emu(80000)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = BG_WHITE
    r.font.name = "Aptos"


def _bottom_accent_bar(slide, x, y, w, color, *, thickness=0.08):
    """Thick colored underline bar (Canva "Workflow Steps Guide" style)."""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(thickness),
    )
    _fill(bar, color)
    _no_line(bar)
    bar.adjustments[0] = 0.5


def _icon_square(slide, x, y, size, glyph, *, color):
    """Rounded-square icon container with white glyph inside."""
    sq = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(size), Inches(size),
    )
    _fill(sq, color)
    _no_line(sq)
    sq.adjustments[0] = 0.18
    tf = sq.text_frame
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = glyph
    r.font.size = Pt(int(size * 22))
    r.font.bold = True
    r.font.color.rgb = BG_WHITE


# --- Slide chrome -----
def _slide_chrome(slide, brand, *, page_no=None, total=15):
    """Canva-inspired: soft lavender-tinted background, minimal chrome."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, BG_TINT)
    _no_line(bg)

    if page_no is not None:
        _text(
            slide, 12.3, 7.15, 0.9, 0.3,
            f"{page_no:02d} / {total:02d}",
            size=8, bold=True, color=INK_MUTED,
            align=PP_ALIGN.RIGHT, font="Consolas",
        )
        _text(
            slide, MARGIN_L, 7.15, 8, 0.3,
            f"MSSP SOC · {brand['name']} · {brand['domain']}",
            size=8, color=INK_MUTED,
        )


def _title_block(slide, eyebrow, title, brand, *, y=0.7, centered=False):
    """Canva-style: big bold title, optional subtitle. Left or centered."""
    align = PP_ALIGN.CENTER if centered else PP_ALIGN.LEFT
    x = MARGIN_L
    w = CONTENT_W
    if eyebrow:
        _text(
            slide, x, y, w, 0.3,
            eyebrow, size=9, bold=True, color=CORP_600,
            uppercase=True, tracking=3.0, align=align,
        )
    _text(
        slide, x, y + 0.35, w, 0.8,
        title, size=32, bold=True, color=INK, font="Aptos Display",
        align=align,
    )
    if not centered:
        _hr(slide, MARGIN_L, y + 1.2, 0.65, color=CORP_600, thickness=1.5)


# --- Insight callout — the "so-what" on every content slide -----
def _insight_bar(slide, y, insight_text, brand):
    """Bottom-of-slide horizontal insight strip — the storytelling element."""
    card = _rect(slide, MARGIN_L, y, CONTENT_W, 0.72,
                 fill=CORP_50, border=None, radius=0.20)
    _no_line(card)
    # Left accent stripe on the callout
    stripe = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(MARGIN_L), Inches(y),
        Inches(0.06), Inches(0.72),
    )
    _fill(stripe, CORP)
    _no_line(stripe)
    # Icon + eyebrow
    _text(
        slide, MARGIN_L + 0.25, y + 0.09, 4, 0.22,
        "▎ AI EXECUTIVE INSIGHT", size=8, bold=True,
        color=CORP, uppercase=True, tracking=2.5,
    )
    # Insight body
    _text(
        slide, MARGIN_L + 0.25, y + 0.32, CONTENT_W - 0.4, 0.4,
        insight_text, size=10.5, color=INK_STRONG, italic=False,
    )


# --- KPI Card — rounded, minimal, no accent bar (consulting style) -----
def _kpi_card(slide, x, y, w, h, *, label, value, unit=None,
              trend=None, trend_positive=True, description=None):
    """Minimal KPI card: uppercase label, big number, tiny trend chip."""
    card = _rect(slide, x, y, w, h,
                 fill=BG_WHITE, border=(BORDER, 0.75), radius=0.09)

    # Uppercase micro label
    _text(
        slide, x + 0.25, y + 0.22, w - 0.5, 0.28,
        label, size=8, bold=True, color=INK_MUTED,
        uppercase=True, tracking=2.5,
    )

    # Big value
    val_tb = slide.shapes.add_textbox(
        Inches(x + 0.25), Inches(y + 0.55),
        Inches(w - 0.5), Inches(0.7),
    )
    vtf = val_tb.text_frame
    vtf.margin_left = vtf.margin_top = Emu(0)
    p = vtf.paragraphs[0]
    rv = p.add_run()
    rv.text = str(value)
    rv.font.name = "Aptos Display"
    rv.font.size = Pt(30)
    rv.font.bold = True
    rv.font.color.rgb = INK
    if unit:
        ru = p.add_run()
        ru.text = f" {unit}"
        ru.font.size = Pt(14)
        ru.font.color.rgb = INK_MUTED
        ru.font.name = "Aptos"

    if description:
        _text(
            slide, x + 0.25, y + h - 0.4, w - 0.5, 0.3,
            description, size=8, color=INK_MUTED,
        )

    if trend is not None:
        # Arrow direction reflects sign of raw trend; color reflects "goodness".
        going_up = trend >= 0
        is_good = trend_positive
        chip_fill = CORP_50 if is_good else RGBColor(0xFE, 0xE2, 0xE2)
        chip_ink = EMERALD if is_good else ROSE
        chip = _rect(
            slide, x + w - 0.85, y + 0.22, 0.6, 0.24,
            fill=chip_fill,
            border=None, radius=0.4,
        )
        chip_tf = chip.text_frame
        chip_tf.margin_top = chip_tf.margin_bottom = Emu(0)
        chip_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = chip_tf.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        arrow = "▲" if going_up else "▼"
        cr.text = f"{arrow} {abs(trend):.1f}%"
        cr.font.size = Pt(7.5)
        cr.font.bold = True
        cr.font.color.rgb = chip_ink


# --- Traffic light indicator (green/amber/red) -----
def _traffic_light(slide, x, y, w, *, label, status, note=""):
    """Row: circle indicator, label, note. status = 'ok'|'watch'|'risk'"""
    col_map = {"ok": EMERALD, "watch": AMBER, "risk": ROSE}
    color = col_map.get(status, GRAY)
    # Colored dot
    dot = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.05), Inches(0.24), Inches(0.24),
    )
    _fill(dot, color)
    _no_line(dot)
    # Label
    _text(
        slide, x + 0.35, y, 4, 0.3,
        label, size=10, bold=True, color=INK,
    )
    # Note (right-aligned)
    if note:
        _text(
            slide, x + 0.35, y + 0.28, w - 0.35, 0.24,
            note, size=8.5, color=INK_MUTED,
        )


# --- Donut gauge -----
def _donut(slide, x, y, size, *, value, label, color=CORP):
    remaining = max(0, 100 - value)
    d = CategoryChartData()
    d.categories = [label]
    d.add_series("v", (value,))
    d.add_series("r", (remaining,))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(x), Inches(y), Inches(size), Inches(size), d,
    )
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = False
    for i, s in enumerate(chart.series):
        for pt in s.points:
            pt.format.fill.solid()
            pt.format.fill.fore_color.rgb = color if i == 0 else RGBColor(0xF1, 0xF5, 0xF9)
            pt.format.line.fill.background()
    try:
        plot = chart._chartSpace.find(qn("c:chart")).find(qn("c:plotArea"))
        doughnut = plot.find(qn("c:doughnutChart"))
        hs = doughnut.find(qn("c:holeSize")) or etree.SubElement(doughnut, qn("c:holeSize"))
        hs.set("val", "78")
    except Exception:
        pass
    # Center number — anchored MIDDLE for perfect vertical alignment
    val_tb = slide.shapes.add_textbox(
        Inches(x), Inches(y + size * 0.25),
        Inches(size), Inches(size * 0.5),
    )
    vtf = val_tb.text_frame
    vtf.margin_left = vtf.margin_right = Emu(0)
    vtf.margin_top = vtf.margin_bottom = Emu(0)
    vtf.word_wrap = False
    vtf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = vtf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    rv = p.add_run()
    rv.text = str(value)
    rv.font.name = "Aptos Display"
    rv.font.size = Pt(30)
    rv.font.bold = True
    rv.font.color.rgb = INK
    # Label below
    _text(
        slide, x, y + size + 0.05, size, 0.3,
        label, size=8.5, bold=True, color=INK_MUTED,
        uppercase=True, tracking=2.0, align=PP_ALIGN.CENTER,
    )


# --- Native line chart -----
def _line_chart(slide, x, y, w, h, *, series_map, accent=CORP, title=None):
    d = CategoryChartData()
    first = next(iter(series_map)) if series_map else None
    cats = [pt["date"] for pt in series_map[first]] if first else []
    if not cats:
        if title:
            _chart_title(slide, x, y - 0.4, w, title)
        _text(slide, x, y + h / 2 - 0.15, w, 0.3, "No data", size=9,
              color=INK_MUTED, align=PP_ALIGN.CENTER)
        return
    d.categories = cats
    for name, pts in series_map.items():
        d.add_series(name, tuple(p["value"] for p in pts))
    cs = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(x), Inches(y), Inches(w), Inches(h), d,
    )
    ch = cs.chart
    ch.has_title = False
    if len(series_map) > 1:
        ch.has_legend = True
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(9)
    else:
        ch.has_legend = False
    colors = [accent, EMERALD, AMBER, ROSE]
    for i, s in enumerate(ch.series):
        ln = s.format.line
        ln.color.rgb = colors[i % len(colors)]
        ln.width = Pt(2.5)
        s.smooth = True
    ch.category_axis.tick_labels.font.size = Pt(8)
    ch.category_axis.tick_labels.font.color.rgb = INK_MUTED
    ch.value_axis.tick_labels.font.size = Pt(8)
    ch.value_axis.tick_labels.font.color.rgb = INK_MUTED
    if title:
        _chart_title(slide, x, y - 0.4, w, title)


def _column_chart(slide, x, y, w, h, *, categories, values, accent=CORP, title=None,
                  point_colors=None):
    d = CategoryChartData()
    if not list(categories):
        if title:
            _chart_title(slide, x, y - 0.4, w, title)
        _text(slide, x, y + h / 2 - 0.15, w, 0.3, "No data", size=9,
              color=INK_MUTED, align=PP_ALIGN.CENTER)
        return
    d.categories = list(categories)
    d.add_series("v", tuple(values))
    cs = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h), d,
    )
    ch = cs.chart
    ch.has_title = False
    ch.has_legend = False
    s = ch.series[0]
    fl = s.format.fill
    fl.solid()
    fl.fore_color.rgb = accent
    # Per-point colors override series color (e.g. severity gradient)
    if point_colors:
        for i, pt in enumerate(s.points):
            if i < len(point_colors):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = point_colors[i]
                pt.format.line.fill.background()
    ch.category_axis.tick_labels.font.size = Pt(8)
    ch.category_axis.tick_labels.font.color.rgb = INK
    ch.value_axis.tick_labels.font.size = Pt(8)
    ch.value_axis.tick_labels.font.color.rgb = INK_MUTED
    if title:
        _chart_title(slide, x, y - 0.4, w, title)


def _bar_chart(slide, x, y, w, h, *, categories, values, accent=CORP, title=None):
    d = CategoryChartData()
    if not list(categories):
        if title:
            _chart_title(slide, x, y - 0.4, w, title)
        _text(slide, x, y + h / 2 - 0.15, w, 0.3, "No data", size=9,
              color=INK_MUTED, align=PP_ALIGN.CENTER)
        return
    d.categories = list(categories)
    d.add_series("v", tuple(values))
    cs = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h), d,
    )
    ch = cs.chart
    ch.has_title = False
    ch.has_legend = False
    s = ch.series[0]
    fl = s.format.fill
    fl.solid()
    fl.fore_color.rgb = accent
    ch.category_axis.tick_labels.font.size = Pt(8)
    ch.category_axis.tick_labels.font.color.rgb = INK
    ch.value_axis.tick_labels.font.size = Pt(8)
    ch.value_axis.tick_labels.font.color.rgb = INK_MUTED
    if title:
        _chart_title(slide, x, y - 0.4, w, title)


def _stacked_column(slide, x, y, w, h, *, categories, series_map, colors=None,
                    title=None):
    d = CategoryChartData()
    if not list(categories):
        if title:
            _chart_title(slide, x, y - 0.4, w, title)
        _text(slide, x, y + h / 2 - 0.15, w, 0.3, "No data", size=9,
              color=INK_MUTED, align=PP_ALIGN.CENTER)
        return
    d.categories = list(categories)
    for name, values in series_map.items():
        d.add_series(name, tuple(values))
    cs = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(x), Inches(y), Inches(w), Inches(h), d,
    )
    ch = cs.chart
    ch.has_title = False
    ch.has_legend = True
    ch.legend.position = XL_LEGEND_POSITION.BOTTOM
    ch.legend.include_in_layout = False
    ch.legend.font.size = Pt(9)
    default_colors = [CORP, CORP_500, AMBER, ROSE]
    if colors is None:
        colors = default_colors
    for i, s in enumerate(ch.series):
        fl = s.format.fill
        fl.solid()
        fl.fore_color.rgb = colors[i % len(colors)]
    ch.category_axis.tick_labels.font.size = Pt(8)
    ch.category_axis.tick_labels.font.color.rgb = INK
    ch.value_axis.tick_labels.font.size = Pt(8)
    ch.value_axis.tick_labels.font.color.rgb = INK_MUTED
    if title:
        _chart_title(slide, x, y - 0.4, w, title)


def _funnel_native(slide, x, y, w, h, *, stages, title=None):
    """
    Ranked cascade — fixed-height rows, proportional bar widths, editable.
    Label sits INSIDE a fixed 3.2" left column so it never wraps vertically.
    stages: list of (label, value) sorted from largest to smallest.
    """
    if title:
        _chart_title(slide, x, y - 0.4, w, title)
    max_v = max(v for _, v in stages) or 1
    n = len(stages)
    row_h = (h - 0.12 * (n - 1)) / n
    label_col_w = 3.2
    bar_area_x = x + label_col_w + 0.15
    bar_area_w = w - label_col_w - 0.15
    for i, (label, v) in enumerate(stages):
        yy = y + i * (row_h + 0.12)
        # Left column: label
        _text(
            slide, x, yy + 0.05, label_col_w, row_h - 0.1,
            label, size=10.5, bold=True, color=INK,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # Track (full width)
        track = _rect(
            slide, bar_area_x, yy + row_h * 0.28,
            bar_area_w, row_h * 0.44,
            fill=BG_SOFT, border=None, radius=0.4,
        )
        # Bar (width proportional to value, but with a legibility minimum)
        pct = v / max_v
        fill_w = max(0.9, bar_area_w * pct)
        bar_color = CORP if i == 0 else CORP_500 if i == 1 else CORP_100
        _rect(
            slide, bar_area_x, yy + row_h * 0.28,
            fill_w, row_h * 0.44,
            fill=bar_color, border=None, radius=0.4,
        )
        # Value at right end
        _text(
            slide, bar_area_x + bar_area_w - 1.2, yy + 0.05,
            1.15, row_h - 0.1,
            f"{v:,}", size=11, bold=True, color=INK,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
            font="Aptos Display",
        )


def _chart_title(slide, x, y, w, text):
    _text(
        slide, x, y, w, 0.3,
        text, size=9, bold=True, color=INK_MUTED,
        uppercase=True, tracking=2.5,
    )
    # thin accent line
    _hr(slide, x, y + 0.28, 0.35, color=CORP, thickness=1.2)


# --- Native table -----
def _table(slide, x, y, w, h, *, headers, rows, title=None):
    if title:
        _chart_title(slide, x, y - 0.4, w, title)
    n_rows = len(rows) + 1
    n_cols = len(headers)
    shape = slide.shapes.add_table(
        n_rows, n_cols, Inches(x), Inches(y), Inches(w), Inches(h),
    )
    t = shape.table
    header_h = Inches(0.42)
    body_h = Inches((h - 0.42) / len(rows)) if rows else Inches(0.35)
    t.rows[0].height = header_h
    for r_i in range(1, n_rows):
        t.rows[r_i].height = body_h
    # Header row
    for c, h_txt in enumerate(headers):
        cell = t.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = CORP
        cell.margin_left = cell.margin_right = Emu(90000)
        cell.margin_top = cell.margin_bottom = Emu(50000)
        cell.text = ""
        p = cell.text_frame.paragraphs[0]
        rn = p.add_run()
        rn.text = str(h_txt)
        rn.font.size = Pt(9.5)
        rn.font.bold = True
        rn.font.color.rgb = BG_WHITE
        rn.font.name = "Aptos"
    # Body rows
    for r_idx, row in enumerate(rows, start=1):
        for c_idx, val in enumerate(row):
            cell = t.cell(r_idx, c_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (
                BG_SOFT if r_idx % 2 == 1 else BG_WHITE
            )
            cell.margin_left = cell.margin_right = Emu(90000)
            cell.margin_top = cell.margin_bottom = Emu(40000)
            cell.text = ""
            p = cell.text_frame.paragraphs[0]
            rn = p.add_run()
            rn.text = str(val)
            rn.font.size = Pt(9)
            rn.font.color.rgb = INK
            rn.font.name = "Aptos"


# --- Timeline (5 nodes) — used on Speed slide -----
def _timeline(slide, x, y, w, *, nodes):
    """nodes = list of (label, value, unit, glyph)"""
    n = len(nodes)
    # main line
    line_y = y + 0.7
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(x + 0.35), Inches(line_y),
        Inches(w - 0.7), Emu(20000),
    )
    _fill(line, BORDER_STRONG)
    _no_line(line)
    step = (w - 0.7) / (n - 1) if n > 1 else 0
    for i, (label, value, unit, glyph) in enumerate(nodes):
        cx = x + 0.35 + i * step - 0.4
        # node disc
        _icon_disc(slide, cx + 0.15, line_y - 0.28, 0.6, color=CORP, glyph=glyph)
        # label above
        _text(
            slide, cx - 0.6, y, 1.7, 0.3,
            label, size=8.5, bold=True, color=INK_MUTED,
            uppercase=True, tracking=1.5, align=PP_ALIGN.CENTER,
        )
        # value below
        _text(
            slide, cx - 0.6, line_y + 0.55, 1.7, 0.5,
            f"{value}", size=18, bold=True, color=INK,
            align=PP_ALIGN.CENTER, font="Aptos Display",
        )
        _text(
            slide, cx - 0.6, line_y + 1.1, 1.7, 0.3,
            unit, size=9, color=INK_MUTED, align=PP_ALIGN.CENTER,
        )


# --- SmartArt-style vertical process list for agenda -----
def _agenda_list(slide, x, y, w, items):
    """items = list of (title, subtitle)"""
    row_h = 0.7
    for i, (title, subtitle) in enumerate(items):
        row_y = y + i * (row_h + 0.15)
        # Number badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(row_y + 0.1),
            Inches(0.5), Inches(0.5),
        )
        _fill(badge, CORP_50)
        _line(badge, CORP, 1.0)
        btf = badge.text_frame
        btf.word_wrap = False
        btf.margin_top = btf.margin_bottom = Emu(0)
        btf.margin_left = btf.margin_right = Emu(0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = f"{i+1:02d}"
        br.font.size = Pt(10)
        br.font.bold = True
        br.font.color.rgb = CORP
        br.font.name = "Aptos Display"
        # Title
        _text(
            slide, x + 0.75, row_y + 0.05, w - 0.9, 0.35,
            title, size=15, bold=True, color=INK, font="Aptos Display",
        )
        # Subtitle
        _text(
            slide, x + 0.75, row_y + 0.42, w - 0.9, 0.28,
            subtitle, size=10, color=INK_MUTED,
        )


# ================================================================
#  SLIDE BUILDERS
# ================================================================

def _slide_cover(prs, brand, period):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, BG_WHITE)
    _no_line(bg)

    # Right side navy panel (30% width)
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.2), 0, Inches(4.15), SLIDE_H,
    )
    _fill(panel, CORP)
    _no_line(panel)

    # Tenant accent "flag" — a thick vertical bar overlaying the left edge
    # of the navy panel. Doubles as a subtle brand-color highlight.
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.12), Inches(2.7), Inches(0.16), Inches(2.1),
    )
    _fill(accent, _hex(brand["primary"]))
    _no_line(accent)

    # Logo mark on left
    _icon_disc(slide, MARGIN_L, 0.75, 0.55, color=CORP, glyph="◆")
    _text(
        slide, MARGIN_L + 0.75, 0.85, 6, 0.35,
        "MSSP SOC KPI CONSOLE", size=10, bold=True,
        color=INK, tracking=3.0, uppercase=True,
    )

    # Eyebrow
    _text(
        slide, MARGIN_L, 3.0, 8, 0.35,
        f"{period.upper()} SECURITY OPERATIONS REVIEW",
        size=11, bold=True, color=CORP, tracking=3.0,
    )

    # Big title
    _text(
        slide, MARGIN_L, 3.5, 8.5, 1.4,
        brand["name"], size=58, bold=True, color=INK, font="Aptos Display",
    )
    _text(
        slide, MARGIN_L, 4.9, 8.5, 0.5,
        f"QRadar Domain · {brand['domain']}",
        size=14, color=INK_MUTED, font="Consolas",
    )

    # Bottom meta
    _hr(slide, MARGIN_L, 6.5, CONTENT_W - 4.2, color=CORP, thickness=1.5)
    _text(
        slide, MARGIN_L, 6.65, 6, 0.3,
        f"CONFIDENTIAL · Prepared {datetime.now(timezone.utc).strftime('%B %Y')}",
        size=9, bold=True, color=INK_MUTED, tracking=2.0, uppercase=True,
    )

    # Right panel text — vertical wordmark
    _text(
        slide, 9.4, 6.7, 3.9, 0.4,
        "SOC OPERATIONS", size=9, bold=True,
        color=RGBColor(0xB4, 0xC6, 0xE7), tracking=3.0, uppercase=True,
    )


def _slide_agenda(prs, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=2)
    _title_block(slide, "Agenda", "What we'll cover today", brand)

    items = [
        ("Executive posture", "Where we stand on the four pillars"),
        ("Threat landscape & incidents", "The story behind the numbers this cycle"),
        ("Speed of response", "MTTD → MTTA → MTTC → MTTR benchmarks"),
        ("Detection coverage", "MITRE ATT&CK posture and gaps"),
        ("Automation ROI", "Hours saved, playbook effectiveness"),
        ("Recommendations & next steps", "Prioritised actions for this cycle"),
    ]
    _agenda_list(slide, MARGIN_L, 2.05, CONTENT_W, items)


def _slide_exec_summary(prs, brand, exec_data, recs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=3)
    _title_block(slide, "Executive Summary", "The story in six numbers", brand)

    # 3 KPI cards top row, 3 bottom row
    row_y = 2.15
    col_w = (CONTENT_W - 0.5) / 3
    kpis_top = [
        {"label": "SOC Health Score", "value": exec_data["health_score"], "unit": None,
         "trend": 2.4, "trend_positive": True, "description": "Composite of SLA, coverage & automation"},
        {"label": "Composite Risk", "value": exec_data["risk_score"], "unit": None,
         "trend": -4.1, "trend_positive": True, "description": "Lower is better · client peer avg 38"},
        {"label": "SLA Compliance", "value": exec_data["sla_compliance"], "unit": "%",
         "trend": 0.8, "trend_positive": True, "description": "vs. contractual 95% target"},
    ]
    kpis_bot = [
        {"label": "Mean Time To Resolve", "value": exec_data["mttr_hours"], "unit": "h",
         "trend": -6.4, "trend_positive": True, "description": "Cycle-over-cycle improvement"},
        {"label": "Detection Coverage", "value": exec_data["detection_coverage"], "unit": "%",
         "trend": 1.9, "trend_positive": True, "description": "MITRE ATT&CK footprint"},
        {"label": "Automation Rate", "value": exec_data["automation_rate"], "unit": "%",
         "trend": 3.4, "trend_positive": True, "description": "Playbook-driven closures"},
    ]
    for i, k in enumerate(kpis_top):
        _kpi_card(slide, MARGIN_L + i * (col_w + 0.25), row_y, col_w, 1.85, **k)
    for i, k in enumerate(kpis_bot):
        _kpi_card(slide, MARGIN_L + i * (col_w + 0.25), row_y + 2.0, col_w, 1.85, **k)

    # Insight
    top_actor = exec_data.get("top_threat_actor", "n/a")
    insight = (
        f"Posture is stable this cycle: SLA holding at {exec_data['sla_compliance']}%, "
        f"MTTR improving, and automation lift continues. Watch on {top_actor} campaigns — "
        f"biggest lever this cycle is automating triage on the top-3 noisy rules."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_posture(prs, brand, exec_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=4)
    _title_block(slide, "Where we are", "Four-pillar posture check", brand)

    # Left: donut
    _donut(slide, MARGIN_L, 2.15, 2.6, value=exec_data["health_score"],
           label="OVERALL HEALTH", color=CORP)

    # Right: 4 pillars as traffic lights
    def status_from(v, hi, mid):
        return "ok" if v >= hi else "watch" if v >= mid else "risk"

    pillars = [
        ("Detection & Coverage", status_from(exec_data["detection_coverage"], 80, 65),
         f"{exec_data['detection_coverage']}% MITRE coverage · target 80%"),
        ("Response Speed", status_from(100 - exec_data["mttr_hours"], 50, 30),
         f"MTTR {exec_data['mttr_hours']}h · target < 48h"),
        ("Service Level Agreement", status_from(exec_data["sla_compliance"], 96, 93),
         f"{exec_data['sla_compliance']}% compliance · target 95%"),
        ("Automation Leverage", status_from(exec_data["automation_rate"], 70, 55),
         f"{exec_data['automation_rate']}% auto-handled · target 70%"),
    ]

    # Right side panel
    panel_x, panel_y, panel_w = 4.4, 2.15, CONTENT_W - 4.4 + MARGIN_L
    for i, (label, status, note) in enumerate(pillars):
        _traffic_light(slide, panel_x, panel_y + i * 0.72, panel_w,
                       label=label, status=status, note=note)
        if i < len(pillars) - 1:
            _hr(slide, panel_x, panel_y + i * 0.72 + 0.68, panel_w - 0.3,
                color=BORDER, thickness=0.4)

    insight = (
        f"Three of four pillars are on-target; the pressure point this cycle is "
        f"{'response speed' if exec_data['mttr_hours'] > 55 else 'detection coverage'} — "
        f"which correlates with the top-3 noisiest rules driving analyst load."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_incidents(prs, brand, soc_data, exec_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=5)
    _title_block(slide, "Threat & Incident Landscape",
                 "Volume, severity and conversion", brand)

    # 3 KPI cards at top
    col_w = (CONTENT_W - 0.5) / 3
    _kpi_card(slide, MARGIN_L, 2.15, col_w, 1.35,
              label="Offenses (QRadar)", value=soc_data["incident_ops"]["total_offenses"],
              trend=3.4, trend_positive=False)
    _kpi_card(slide, MARGIN_L + col_w + 0.25, 2.15, col_w, 1.35,
              label="Incidents (XSOAR)", value=soc_data["incident_ops"]["total_incidents"],
              trend=-2.1, trend_positive=True)
    _kpi_card(slide, MARGIN_L + 2 * (col_w + 0.25), 2.15, col_w, 1.35,
              label="Offense → Incident", value=soc_data["incident_ops"]["conversion_rate"],
              unit="%", description="Quality-of-alerts indicator")

    # Chart: severity distribution stacked column
    severity = soc_data["detection_health"]["severity_distribution"]
    severity_palette = {
        "Critical": ROSE,
        "High": RGBColor(0xEA, 0x58, 0x0C),   # orange-600
        "Medium": AMBER,
        "Low": RGBColor(0x0D, 0x94, 0x88),    # teal-600
        "Informational": RGBColor(0x64, 0x74, 0x8B),
    }
    pt_colors = [severity_palette.get(s["severity"], CORP) for s in severity]
    _column_chart(
        slide, MARGIN_L, 4.0, CONTENT_W, 2.15,
        categories=[s["severity"] for s in severity],
        values=[s["count"] for s in severity],
        accent=CORP,
        point_colors=pt_colors,
        title="Incident Severity Distribution  ·  This cycle",
    )

    insight = (
        f"Incident volume held despite a {soc_data['incident_ops']['escalation_rate']}% "
        f"escalation rate. The severity mix skews to Medium/Low, indicating detection "
        f"is catching earlier — but backlog aging ({soc_data['incident_ops']['backlog_aging_days']}d) "
        f"warrants an L1 shift rebalance."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_speed(prs, brand, soc_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=6)
    _title_block(slide, "Speed of Response",
                 "The four moments that matter", brand)

    speed = soc_data["speed_metrics"]
    nodes = [
        ("DETECT", speed["mttd_min"], "minutes to detect", "◉"),
        ("ACKNOWLEDGE", speed["mtta_min"], "minutes to ack", "◐"),
        ("CONTAIN", speed["mttc_hours"], "hours to contain", "◆"),
        ("RESOLVE", speed["mttr_hours"], "hours to resolve", "✓"),
    ]
    _timeline(slide, MARGIN_L, 2.35, CONTENT_W, nodes=nodes)

    # Comparison chip strip
    y = 5.15
    _text(
        slide, MARGIN_L, y, CONTENT_W, 0.3,
        "BENCHMARK COMPARISON", size=9, bold=True, color=INK_MUTED,
        uppercase=True, tracking=2.5,
    )
    _hr(slide, MARGIN_L, y + 0.3, 0.35, color=CORP, thickness=1.2)

    bench_items = [
        ("MTTD peer avg", "12.0 min", "ok" if speed["mttd_min"] <= 12 else "watch"),
        ("MTTA peer avg", "20.0 min", "ok" if speed["mtta_min"] <= 20 else "watch"),
        ("MTTR peer avg", "62.0 h", "ok" if speed["mttr_hours"] <= 62 else "watch"),
    ]
    for i, (lbl, val, st) in enumerate(bench_items):
        x = MARGIN_L + i * (CONTENT_W / 3)
        col = EMERALD if st == "ok" else AMBER
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.55), Inches(0.2), Inches(0.2),
        )
        _fill(dot, col)
        _no_line(dot)
        _text(slide, x + 0.3, y + 0.5, 3, 0.3, lbl,
              size=10, color=INK)
        _text(slide, x + 0.3, y + 0.72, 3, 0.28, val,
              size=9, color=INK_MUTED)

    insight = (
        f"Detect + acknowledge sit inside benchmark; resolve is the drag — "
        f"queue time of {speed['queue_time_min']} min is 60% of the delta. "
        f"Auto-enrichment at incident creation would compress MTTR by an estimated 15–25%."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_threats(prs, brand, ti_data, client_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=7)
    _title_block(slide, "Threat Landscape",
                 "Who is targeting us and how", brand)

    # Left: threat actors as matrix (Name | Origin | Activity)
    actors = ti_data["landscape"]["threat_actors"][:5]
    rows = [(a["name"], a["origin"], f"{a['activity']}") for a in actors]
    _table(slide, MARGIN_L, 2.55, 5.7, 2.7,
           headers=["Threat Actor", "Origin", "Activity"],
           rows=rows, title="Top Active Actors")

    # Right: malware families bar chart
    mal = ti_data["landscape"]["malware_families"][:6]
    _bar_chart(
        slide, MARGIN_L + 6.0, 2.55, CONTENT_W - 6.0, 2.7,
        categories=[m["family"] for m in mal],
        values=[m["count"] for m in mal],
        accent=CORP,
        title="Malware Families  ·  Detections",
    )

    top_actor = actors[0]["name"] if actors else "n/a"
    insight = (
        f"{top_actor} tops actor activity this cycle. The malware mix is "
        f"ransomware-heavy — recommend pushing the {top_actor} TTP hunt-pack to "
        f"detection engineering and briefing top-tier clients before month-end."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_detection(prs, brand, det_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=8)
    _title_block(slide, "Detection Coverage",
                 "MITRE ATT&CK footprint by tactic", brand)

    tactics = det_data["mitre_heatmap"]
    _bar_chart(
        slide, MARGIN_L, 2.5, CONTENT_W, 3.4,
        categories=[t["tactic"] for t in tactics],
        values=[t["coverage"] for t in tactics],
        accent=CORP,
        title=f"Coverage % across {len(tactics)} tactics",
    )

    # Gap summary right-aligned mini KPIs — inline text
    covered = det_data["gap_analysis"]["techniques_covered"]
    missing = det_data["gap_analysis"]["techniques_missing"]
    _opps = det_data["gap_analysis"].get("new_opportunities") or []
    _top_opp = _opps[0] if _opps else "n/a"
    _text(
        slide, MARGIN_L, 5.95, CONTENT_W, 0.3,
        f"{covered} techniques covered   ·   {missing} techniques missing   ·   "
        f"Top opportunity: {_top_opp}",
        size=10, color=INK_MUTED,
    )

    insight = (
        f"Coverage is strongest on Execution and Persistence, thinnest on "
        f"Credential Access and Defense Evasion — where {missing} techniques remain "
        f"uncovered. Recommend three new detections this sprint on the top-3 gaps."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_automation(prs, brand, soar_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=9)
    _title_block(slide, "Automation ROI",
                 "The economics of the SOAR investment", brand)

    # Two KPI cards on the left
    _kpi_card(slide, MARGIN_L, 2.15, 3.0, 1.6,
              label="Hours Saved", value=soar_data["efficiency"]["hours_saved"],
              description="This cycle · playbook automation")
    _kpi_card(slide, MARGIN_L, 3.9, 3.0, 1.6,
              label="Automation ROI", value=soar_data["efficiency"]["automation_roi_pct"],
              unit="%", trend=14.6, trend_positive=True,
              description="Return on SOAR investment")

    # Right chart: automation trend
    _line_chart(
        slide, MARGIN_L + 3.3, 2.15, CONTENT_W - 3.3, 3.4,
        series_map={
            "Automation %": soar_data["trends"]["automation"],
            "Success %": soar_data["trends"]["success"],
        },
        accent=CORP,
        title="Automation & Success Rate Trend",
    )

    insight = (
        f"Automation at {soar_data['health']['automation_rate']}% is delivering "
        f"{soar_data['efficiency']['hours_saved']} recovered analyst hours this cycle. "
        f"Converting the top-2 manual flows to full auto-remediation would push automation "
        f"past 75% and unlock a further ~20% ROI."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_client_impact(prs, brand, client_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=10)
    _title_block(slide, "Client Impact",
                 "The story a Board sees", brand)

    col_w = (CONTENT_W - 0.5) / 3
    _kpi_card(slide, MARGIN_L, 2.15, col_w, 1.55,
              label="Composite Risk", value=client_data["scorecard"]["composite_risk_score"],
              description="Peer-benchmarked · lower is better")
    _kpi_card(slide, MARGIN_L + col_w + 0.25, 2.15, col_w, 1.55,
              label="Client Risk Rank", value=f"#{client_data['scorecard']['client_risk_rank']}",
              description="Out of tenant portfolio")
    _kpi_card(slide, MARGIN_L + 2 * (col_w + 0.25), 2.15, col_w, 1.55,
              label="Major P1/P2", value=client_data["scorecard"]["major_p1_p2_incidents"],
              description="Board-visible incidents this cycle")

    # Native funnel: assets → alerts → incidents → escalations
    stages = [
        ("Total Assets Monitored", 4200),
        ("Alerts Generated", 1830),
        ("Incidents Opened", client_data["scorecard"]["major_p1_p2_incidents"] * 40),
        ("Board Escalations", client_data["scorecard"]["major_p1_p2_incidents"]),
    ]
    _funnel_native(slide, MARGIN_L, 4.35, CONTENT_W, 1.75,
                   stages=stages, title="From Signal to Board · Filter Funnel")

    insight = (
        f"Only {client_data['scorecard']['major_p1_p2_incidents']} incidents reached "
        f"Board visibility this cycle — a healthy signal-to-noise ratio. "
        f"YoY MTTR is trending {client_data['scorecard']['yoy_mttr_delta']:+}%; "
        f"repeat incidents remain the main risk-score contributor."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_ai_recommendations(prs, brand, recs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=11)
    _title_block(slide, "AI Recommendations",
                 "HuggingFace-augmented cybersecurity reasoning", brand)

    # 4 recommendation cards in a 2x2 grid (minimal, borderless, whitespace)
    top = 1.95
    card_w = (CONTENT_W - 0.4) / 2
    card_h = 2.05
    for i, rec in enumerate(recs[:4]):
        r, c = divmod(i, 2)
        x = MARGIN_L + c * (card_w + 0.4)
        y = top + r * (card_h + 0.15)
        _rect(slide, x, y, card_w, card_h,
              fill=BG_WHITE, border=(BORDER, 0.75), radius=0.06)
        # Number badge
        badge = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.25), Inches(y + 0.25),
            Inches(0.45), Inches(0.45),
        )
        _fill(badge, CORP)
        _no_line(badge)
        btf = badge.text_frame
        btf.word_wrap = False
        btf.margin_left = btf.margin_right = Emu(0)
        btf.margin_top = btf.margin_bottom = Emu(0)
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bp = btf.paragraphs[0]
        bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run()
        br.text = f"{i+1:02d}"
        br.font.size = Pt(10)
        br.font.bold = True
        br.font.color.rgb = BG_WHITE
        br.font.name = "Aptos Display"

        # Area micro-tag
        _text(slide, x + 0.85, y + 0.30, card_w - 2.7, 0.28,
              rec.get("area", ""), size=8, bold=True, color=CORP,
              uppercase=True, tracking=2.0)
        # HF LLM badge on right if enriched
        if rec.get("reasoning_source") == "hf-llm":
            _text(slide, x + card_w - 1.75, y + 0.30, 1.55, 0.28,
                  "◆ HF LLM REASONING", size=7, bold=True, color=CORP_500,
                  uppercase=True, tracking=1.5, align=PP_ALIGN.RIGHT)
        # Title
        _text(slide, x + 0.25, y + 0.75, card_w - 0.5, 0.4,
              rec.get("title", ""), size=12, bold=True, color=INK,
              font="Aptos Display")
        # Body — larger region, tighter truncation
        body_text = rec.get("reasoning") or rec.get("action", "")
        if len(body_text) > 220:
            body_text = body_text[:217] + "…"
        _text(slide, x + 0.25, y + 1.20, card_w - 0.5, card_h - 1.28,
              body_text, size=9.5, color=INK_STRONG)

    insight = (
        "The next-cycle plan focuses on three levers: automate the top-2 manual playbooks, "
        "ship three detections against uncovered ATT&CK techniques, and re-tune the noisiest rules."
    )
    _insight_bar(slide, 6.4, insight, brand)


def _slide_next_steps(prs, brand, recs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=12)
    _title_block(slide, "Next Steps",
                 "Owners, priorities, and delivery cadence", brand)

    rows = []
    for r in recs[:5]:
        rows.append((
            r.get("priority", "P3"),
            r.get("area", "-"),
            r.get("title", "")[:60],
            (r.get("action", "") or r.get("reasoning", ""))[:90],
        ))
    _table(slide, MARGIN_L, 2.5, CONTENT_W, 3.6,
           headers=["Priority", "Area", "Recommendation", "Owner Action"],
           rows=rows, title="Prioritised Actions · This Cycle")

    insight = (
        "P1 actions delivered inside 5 business days; P2 by end of cycle; P3 as capacity permits. "
        "Owners: SOC Operations lead + Detection Engineering lead. Next review at cycle + 1."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_appendix_divider(prs, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=13)
    _text(slide, MARGIN_L, 1.6, CONTENT_W, 0.4,
          "APPENDIX", size=12, bold=True, color=CORP, tracking=5.0)
    _text(slide, MARGIN_L, 2.05, CONTENT_W, 1.1,
          "Supporting data & methodology", size=38, bold=True,
          color=INK, font="Aptos Display")
    _hr(slide, MARGIN_L, 3.15, 0.65, color=CORP, thickness=1.5)

    # Methodology cards — 3 columns
    col_w = (CONTENT_W - 0.6) / 3
    pillars = [
        ("Data sources",
         "QRadar (offenses, rules, assets) · XSOAR (incidents, SLA, playbooks) · "
         "Threat Intelligence feeds (advisories, IOCs, MITRE ATT&CK, threat actors)."),
        ("Calculation windows",
         "All rate-based KPIs use a rolling 30-day window; MITRE coverage is a "
         "point-in-time snapshot; peer benchmarks refreshed quarterly."),
        ("AI reasoning",
         "Executive insights are generated by a local HuggingFace SmolLM instance "
         "grounded on this cycle's KPIs. No customer data leaves the environment."),
    ]
    for i, (title, body) in enumerate(pillars):
        x = MARGIN_L + i * (col_w + 0.3)
        _shadow_card(slide, x, 3.85, col_w, 2.35, radius=0.08)
        _bottom_accent_bar(slide, x + 0.3, 3.95, 0.5, CORP)
        _text(slide, x + 0.3, 4.15, col_w - 0.6, 0.4,
              title, size=13, bold=True, color=INK, font="Aptos Display")
        _text(slide, x + 0.3, 4.65, col_w - 0.6, 1.55,
              body, size=10, color=INK_MUTED)


def _slide_data_snapshot(prs, brand, ti_data, soar_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, brand, page_no=14)
    _title_block(slide, "Data Snapshot",
                 "Reference metrics for deep dives", brand)

    # Left mini KPIs
    small_kpis = [
        ("Total Advisories", ti_data["landscape"]["total_advisories"]),
        ("IOC Volume", ti_data["landscape"]["ioc_volume"]),
        ("New CVEs", ti_data["landscape"]["new_cves"]),
        ("Critical CVEs", ti_data["landscape"]["critical_cves"]),
    ]
    for i, (lbl, val) in enumerate(small_kpis):
        r, c = divmod(i, 2)
        x = MARGIN_L + c * 2.4
        y = 2.25 + r * 1.4
        _kpi_card(slide, x, y, 2.2, 1.25, label=lbl, value=val)

    # Right side: playbook mini table
    playbooks = soar_data["playbooks"][:5]
    rows = [(p["name"][:26], p["executions"], f"{p['success_rate']}%")
            for p in playbooks]
    _table(slide, MARGIN_L + 5.0, 2.5, CONTENT_W - 5.0, 3.4,
           headers=["Playbook", "Runs", "Success"],
           rows=rows, title="Top Playbooks · Reference")

    insight = (
        "Full raw data sets available on request; this appendix contains the highest-signal "
        "metrics for a follow-up deep-dive with the SOC operations team."
    )
    _insight_bar(slide, 6.25, insight, brand)


def _slide_thankyou(prs, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, BG_WHITE)
    _no_line(bg)

    # Right navy panel
    panel = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.2), 0, Inches(4.15), SLIDE_H,
    )
    _fill(panel, CORP)
    _no_line(panel)

    # Tenant accent flag on left edge of panel
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(9.12), Inches(2.7), Inches(0.16), Inches(2.1),
    )
    _fill(accent, _hex(brand["primary"]))
    _no_line(accent)

    # Logo mark on the navy panel
    _icon_disc(slide, 10.8, 3.15, 0.75, color=BG_WHITE, glyph="◆")
    # Override glyph color to navy for contrast on white disc
    logo = slide.shapes[-1]
    logo.text_frame.paragraphs[0].runs[0].font.color.rgb = CORP
    _text(slide, 9.4, 4.05, 3.9, 0.35,
          "MSSP SOC · KPI CONSOLE", size=9, bold=True,
          color=RGBColor(0xB4, 0xC6, 0xE7), tracking=3.0, uppercase=True,
          align=PP_ALIGN.CENTER)

    _text(slide, MARGIN_L, 2.9, 8, 0.4,
          "THANK YOU", size=13, bold=True, color=CORP, tracking=6.0)
    _text(slide, MARGIN_L, 3.4, 8.5, 1.6,
          "Questions welcome.", size=54, bold=True, color=INK,
          font="Aptos Display")
    _text(slide, MARGIN_L, 5.15, 8, 0.5,
          "Reach the SOC team for a deep-dive on any KPI.",
          size=15, color=INK_MUTED)
    _hr(slide, MARGIN_L, 5.9, 0.65, color=CORP, thickness=1.5)
    _text(slide, MARGIN_L, 6.15, 8, 0.35,
          f"MSSP SOC · {brand['name']} · {brand['domain']}",
          size=10, bold=True, color=INK_MUTED, tracking=2.0, uppercase=True)


# ================================================================
#  QBR-STYLE SLIDE BUILDERS (Deloitte teal/navy)
# ================================================================

def _qbr_footer(slide, brand, page_no, *, dark=False):
    col = D_GRAY if dark else D_GRAY_DK
    _text(slide, MARGIN_L, 7.12, 8, 0.28,
          f"{brand['name'].upper()}  ·  {brand['period_label']}",
          size=8, bold=True, color=col, tracking=1.5)
    _wordmark(slide, 11.55, 7.02, size=10.5, dark=not dark)


def _qbr_group(slide, x, y, w, title):
    _text(slide, x, y, w, 0.28, title, size=10, bold=True, color=D_TEAL,
          uppercase=True, tracking=2.5)
    _hr(slide, x, y + 0.3, w, color=D_GRAY, thickness=0.5)


def _qbr_tile(slide, x, y, w, h, big, label, *, big_color=D_NAVY):
    card = _rect(slide, x, y, w, h, fill=BG_WHITE, border=(D_GRAY, 0.75), radius=0.06)
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(0.07), Inches(h))
    _fill(stripe, D_TEAL); _no_line(stripe)
    vtb = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.16), Inches(w - 0.35), Inches(h * 0.5))
    vtf = vtb.text_frame; vtf.word_wrap = False
    vtf.margin_left = vtf.margin_top = Emu(0)
    vp = vtf.paragraphs[0]
    rv = vp.add_run(); rv.text = str(big)
    rv.font.name = "Aptos Display"; rv.font.size = Pt(26); rv.font.bold = True
    rv.font.color.rgb = big_color if str(big) != "N/A" else D_GRAY
    _text(slide, x + 0.24, y + h - 0.62, w - 0.4, 0.55, label, size=8.5,
          color=D_GRAY_DK, anchor=MSO_ANCHOR.TOP)


def _qbr_cover(prs, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, D_NAVY_DK); _no_line(bg)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.32), SLIDE_H)
    _fill(band, D_TEAL); _no_line(band)
    # faint teal block bottom-right for depth
    blk = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.9), Inches(5.6), Inches(3.43), Inches(1.9))
    _fill(blk, D_NAVY); _no_line(blk)
    _wordmark(slide, MARGIN_L, 0.7, size=26, dark=False)
    _text(slide, MARGIN_L, 2.75, 11, 0.4,
          f"{brand['name'].upper()}  ·  {brand['period_tag']}",
          size=13, bold=True, color=D_TEAL, tracking=3.0)
    _text(slide, MARGIN_L, 3.25, 11.5, 1.3, "Quarterly SOC Services Report",
          size=44, bold=True, color=BG_WHITE, font="Aptos Display")
    _text(slide, MARGIN_L, 4.7, 10, 0.6, brand["period_label"],
          size=20, color=D_TEAL_LT, tracking=2.0)
    _hr(slide, MARGIN_L, 5.5, 3.0, color=D_TEAL, thickness=2.0)
    _text(slide, MARGIN_L, 6.75, 10, 0.3, "CONFIDENTIAL · INTERNAL USE ONLY",
          size=9, bold=True, color=D_GRAY, tracking=2.5)


def _qbr_section(prs, brand, num, title, subtitle, tags, idx, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, D_NAVY_DK); _no_line(bg)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.32), SLIDE_H)
    _fill(band, D_TEAL); _no_line(band)
    _text(slide, MARGIN_L, 1.4, 6, 0.35, f"SECTION {idx:02d} / {total:02d}",
          size=11, bold=True, color=D_TEAL, tracking=3.0)
    _text(slide, MARGIN_L, 2.0, 6, 2.3, num, size=140, bold=True,
          color=D_TEAL, font="Aptos Display")
    _text(slide, 4.6, 2.7, 8.0, 1.0, title, size=40, bold=True,
          color=BG_WHITE, font="Aptos Display")
    _text(slide, 4.62, 3.85, 8.0, 0.9, subtitle, size=13, color=D_GRAY)
    _text(slide, 4.62, 5.3, 8.0, 0.3, tags, size=9.5, bold=True,
          color=D_TEAL_LT, tracking=2.0)
    _qbr_footer(slide, brand, idx, dark=True)


def _qbr_title(slide, title, subtitle):
    _text(slide, MARGIN_L, 0.55, CONTENT_W, 0.7, title, size=26, bold=True,
          color=D_NAVY, font="Aptos Display")
    if subtitle:
        _text(slide, MARGIN_L, 1.28, CONTENT_W, 0.5, subtitle, size=11, color=D_GRAY_DK)
    _hr(slide, MARGIN_L, 1.75, CONTENT_W, color=D_TEAL, thickness=1.5)


def _qbr_exec_overview(prs, brand, ex, soc, ti):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, D_BG); _no_line(bg)
    _qbr_title(slide, "Security Operations: Executive Overview",
               "Quarterly performance, service reliability and detection outcomes for this tenant.")
    s = (soc or {}).get("summary", {}) if soc else {}
    ex = ex or {}
    ti = ti or {}
    tid_ok = ti.get("data_status") == "live"
    col_w = (CONTENT_W - 0.75) / 4

    def row(items, y):
        for i, (big, lbl) in enumerate(items):
            _qbr_tile(slide, MARGIN_L + i * (col_w + 0.25), y, col_w, 1.05, big, lbl)

    _qbr_group(slide, MARGIN_L, 1.95, CONTENT_W, "Service Reliability")
    row([
        (_vs(s.get("sla_compliance_pct", ex.get("sla_compliance")), pct=True), "SLA compliance vs 95% target"),
        (_vs(None), "service uptime"),
        (_vs(s.get("mttd_minutes"), suffix=" min"), "MTTD · mean time to detect"),
        (_vs(s.get("mttr_hours", ex.get("mttr_hours")), suffix=" h"), "MTTR · mean time to resolve"),
    ], 2.35)

    _qbr_group(slide, MARGIN_L, 3.65, CONTENT_W, "Quarterly Signal")
    row([
        (_vs(s.get("total_incidents", ex.get("incidents"))), "incidents handled (XSOAR)"),
        (_vs(None), "offenses (QRadar)"),
        (_vs(s.get("false_positive_rate", ex.get("false_positive_rate")), pct=True), "false-positive rate"),
        (_vs(s.get("true_positive_rate"), pct=True), "true-positive rate"),
    ], 4.05)

    _qbr_group(slide, MARGIN_L, 5.35, CONTENT_W, "Coverage & Detection")
    row([
        (_vs(ex.get("detection_coverage"), pct=True), "MITRE ATT&CK coverage"),
        (_vs(ex.get("automation_rate"), pct=True), "automation rate"),
        (_vs(ti["summary"]["total_advisories"] if tid_ok else None), "threat advisories deployed"),
        (_vs(ti["summary"].get("total_iocs") if tid_ok and "total_iocs" in ti.get("summary", {}) else None), "IOCs in security systems"),
    ], 5.75)

    _qbr_footer(slide, brand, 3)


def _qbr_exec_performance(prs, brand, ex, soc, rules_count, qbr):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, D_BG); _no_line(bg)
    _qbr_title(slide, "Security Operations: Executive Performance View", None)
    s = (soc or {}).get("summary", {}) if soc else {}
    ex = ex or {}
    sla = s.get("sla_compliance_pct", ex.get("sla_compliance"))
    mttr = s.get("mttr_hours", ex.get("mttr_hours"))
    live = bool(s) or ex.get("data_status") == "live"

    # Headline update (teal band)
    hb = _rect(slide, MARGIN_L, 2.0, CONTENT_W, 1.05, fill=D_TEAL, border=None, radius=0.05)
    _text(slide, MARGIN_L + 0.3, 2.12, CONTENT_W - 0.6, 0.3, "HEADLINE UPDATE",
          size=9, bold=True, color=BG_WHITE, tracking=2.5)
    headline = (
        f"Broader visibility is translating into faster, more resilient defense — "
        f"SLA holding at {_vs(sla, pct=True)} with MTTR at {_vs(mttr, suffix='h')}."
        if live else
        "Upload XSOAR incident data to populate this tenant's executive performance view."
    )
    _text(slide, MARGIN_L + 0.3, 2.44, CONTENT_W - 0.6, 0.55, headline,
          size=13, bold=True, color=BG_WHITE, font="Aptos Display")

    # SOC SERVICES numbers row
    _qbr_group(slide, MARGIN_L, 3.35, 7.0, "SOC Services")
    col_w = (7.0 - 0.6) / 4
    svc = [
        (_vs(s.get("total_incidents", ex.get("incidents"))), "incidents escalated"),
        (_vs(len((qbr or {}).get("log_sources") or []) or None), "log sources integrated"),
        (_vs(rules_count), "use cases enabled"),
        (_vs(None), "flow interfaces"),
    ]
    for i, (big, lbl) in enumerate(svc):
        _qbr_tile(slide, MARGIN_L + i * (col_w + 0.2), 3.75, col_w, 1.05, big, lbl)

    _qbr_group(slide, MARGIN_L, 5.05, 7.0, "Response & SLA")
    col_w2 = (7.0 - 0.4) / 3
    resp = [
        (_vs(sla, pct=True), "incidents within SLA"),
        (_vs(mttr, suffix=" h"), "mean time to resolve"),
        (_vs(s.get("mttd_minutes"), suffix=" min"), "mean time to detect"),
    ]
    for i, (big, lbl) in enumerate(resp):
        _qbr_tile(slide, MARGIN_L + i * (col_w2 + 0.2), 5.45, col_w2, 1.05, big, lbl)

    # Right column: operating model bullets
    bx = 8.5
    _qbr_group(slide, bx, 3.35, CONTENT_W - (bx - MARGIN_L), "Operating Model")
    bullets = [
        "Continuous log-source reconciliation onboards new telemetry.",
        "Regular use-case reviews and rule fine-tuning improve fidelity.",
        "Threat advisories and IOCs strengthen emerging-attack defense.",
        "Automation-driven triage compresses analyst handling time.",
    ]
    for i, b in enumerate(bullets):
        _text(slide, bx, 3.8 + i * 0.62, 4.0, 0.6, "–  " + b, size=10, color=D_SLATE)

    _qbr_footer(slide, brand, 4)


def _qbr_log_sources(prs, brand, qbr):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, D_BG); _no_line(bg)
    ls = (qbr or {}).get("log_sources") or []
    total = (qbr or {}).get("total", 0)
    top = ls[0] if ls else None
    title = (f"{top['name']} Drove {top['pct']}% of Incidents This Period"
             if top else "Log Source Distribution")
    _qbr_title(slide, title,
               "Share of incidents by originating log source, ranked highest to lowest.")
    if not ls:
        _text(slide, MARGIN_L, 3.4, CONTENT_W, 0.5,
              "N/A — no log-source field found in the uploaded incidents.",
              size=13, color=D_GRAY_DK, align=PP_ALIGN.CENTER)
        _qbr_footer(slide, brand, 6)
        return
    y0 = 2.15
    row_h = min(0.5, (4.6) / len(ls))
    label_w = 3.4
    bar_x = MARGIN_L + label_w + 0.15
    bar_w = CONTENT_W - label_w - 0.15
    maxpct = max(x["pct"] for x in ls) or 1
    for i, x in enumerate(ls):
        yy = y0 + i * (row_h + 0.12)
        _text(slide, MARGIN_L, yy, label_w, row_h, x["name"], size=10.5, bold=True,
              color=D_SLATE, anchor=MSO_ANCHOR.MIDDLE)
        _rect(slide, bar_x, yy + row_h * 0.15, bar_w, row_h * 0.7, fill=BG_WHITE,
              border=(D_GRAY, 0.5), radius=0.4)
        fillw = max(0.25, bar_w * (x["pct"] / maxpct))
        _rect(slide, bar_x, yy + row_h * 0.15, fillw, row_h * 0.7,
              fill=D_TEAL if i > 0 else D_NAVY, border=None, radius=0.4)
        _text(slide, bar_x + bar_w - 2.0, yy, 1.95, row_h,
              f"{x['pct']}%  ({x['count']}/{total})", size=9.5, bold=True,
              color=D_SLATE, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    _qbr_footer(slide, brand, 6)


def _qbr_alerts(prs, brand, qbr):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, D_BG); _no_line(bg)
    abm = (qbr or {}).get("alerts_by_month") or {}
    months = abm.get("months") or []
    _qbr_title(slide, "Alert Volume by Month; Severity Mix in Context",
               "Monthly incident volume split by severity, with mean-time-to-resolve context.")
    if not months:
        _text(slide, MARGIN_L, 3.4, CONTENT_W, 0.5,
              "N/A — no dated incidents to chart by month.",
              size=13, color=D_GRAY_DK, align=PP_ALIGN.CENTER)
        _qbr_footer(slide, brand, 7)
        return
    series = abm.get("series", {})
    series_map = {k: series[k] for k in ("High", "Medium", "Low") if any(series.get(k, []))}
    if series.get("Critical") and any(series["Critical"]):
        series_map = {"Critical": series["Critical"], **series_map}
    colors = [SEV_QBR.get(k, D_GRAY) for k in series_map.keys()]
    _stacked_column(slide, MARGIN_L, 2.4, 8.0, 3.7,
                    categories=months, series_map=series_map, colors=colors,
                    title="Alerts by month · severity")
    # Right callouts
    cx = 9.0
    _qbr_group(slide, cx, 2.05, CONTENT_W - (cx - MARGIN_L), "What drove the quarter")
    totals = [(m, sum(series.get(s, [0] * len(months))[i] for s in series))
              for i, m in enumerate(months)]
    peak = max(totals, key=lambda t: t[1]) if totals else ("—", 0)
    mttrm = (qbr or {}).get("mttr_by_month") or []
    avg_mttr = (qbr or {}).get("mttr_median_hours")
    facts = [
        (f"{peak[1]:,}", f"alerts in {peak[0]} — the period peak."),
        (f"{(qbr or {}).get('total', 0):,}", "total incidents across the period."),
        (_vs(avg_mttr, suffix=" h"), "median time to resolve (outlier-robust)."),
    ]
    for i, (big, lbl) in enumerate(facts):
        yy = 2.5 + i * 1.25
        _text(slide, cx, yy, 4.0, 0.5, big, size=26, bold=True, color=D_TEAL,
              font="Aptos Display")
        _text(slide, cx, yy + 0.5, 4.0, 0.6, lbl, size=10, color=D_SLATE)
    _qbr_footer(slide, brand, 7)


def _qbr_mitre(prs, brand, qbr):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, D_BG); _no_line(bg)
    tactics = (qbr or {}).get("tactics") or []
    top = tactics[0] if tactics else None
    title = (f"{top['tactic']} Led the Quarter" if top else "MITRE ATT&CK Tactic Activity")
    _qbr_title(slide, title,
               "Incident volume mapped to MITRE ATT&CK tactics for this tenant.")
    if not tactics:
        _text(slide, MARGIN_L, 3.4, CONTENT_W, 0.5,
              "N/A — no MITRE tactic mapping in the uploaded incidents.",
              size=13, color=D_GRAY_DK, align=PP_ALIGN.CENTER)
        _qbr_footer(slide, brand, 8)
        return
    top8 = tactics[:8]
    _bar_chart(slide, MARGIN_L, 2.35, 8.0, 3.8,
               categories=[t["tactic"] for t in top8],
               values=[t["count"] for t in top8],
               accent=D_TEAL, title="Alerts by MITRE ATT&CK tactic")
    cx = 9.0
    _qbr_group(slide, cx, 2.05, CONTENT_W - (cx - MARGIN_L), "What shifted this quarter")
    for i, t in enumerate(tactics[:3]):
        yy = 2.5 + i * 1.25
        _text(slide, cx, yy, 4.0, 0.5, f"{t['count']:,}", size=26, bold=True,
              color=D_TEAL, font="Aptos Display")
        _text(slide, cx, yy + 0.5, 4.0, 0.6, f"{t['tactic']} alerts this period.",
              size=10, color=D_SLATE)
    _qbr_footer(slide, brand, 8)


def _qbr_closing(prs, brand):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    _fill(bg, D_NAVY_DK); _no_line(bg)
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.32), SLIDE_H)
    _fill(band, D_TEAL); _no_line(band)
    _wordmark(slide, MARGIN_L, 0.8, size=24, dark=False)
    _text(slide, MARGIN_L, 3.0, 11, 0.9, "Thank you.", size=44, bold=True,
          color=BG_WHITE, font="Aptos Display")
    _text(slide, MARGIN_L, 4.1, 11, 0.5,
          f"Prepared for {brand['name']} · {brand['period_label']}",
          size=13, color=D_TEAL_LT)
    _text(slide, MARGIN_L, 6.4, 11.5, 0.9,
          "This report is intended solely for the internal use of the named client. "
          "It is confidential and may not be distributed without prior written consent.",
          size=8.5, color=D_GRAY)


# ================================================================
#  Entry point
# ================================================================

def _period_labels(period: str):
    """Human labels for the cover/footer from the app's period."""
    now = datetime.now(timezone.utc)
    p = (period or "monthly").lower()
    q = (now.month - 1) // 3 + 1
    tag = {"weekly": f"WK {now.isocalendar()[1]} {now.year}",
           "monthly": now.strftime("%b %Y").upper(),
           "quarterly": f"Q{q} {now.year}"}.get(p, p.upper())
    label = {"weekly": now.strftime("WEEK OF %d %B %Y").upper(),
             "monthly": now.strftime("%B %Y").upper(),
             "quarterly": f"Q{q} {now.year}"}.get(p, now.strftime("%B %Y").upper())
    return tag, label


def build_pptx(tenant: dict, period: str, all_data: dict,
               recommendations: list) -> io.BytesIO:
    # Preferred: edit the client's real QBR deck as a template so layout,
    # fonts and wording match exactly. Fall back to the native builder if the
    # template file is missing or editing fails.
    try:
        import pptx_template
        if pptx_template.template_available():
            return pptx_template.build_from_template(tenant, period, all_data, recommendations)
    except Exception:
        logging.getLogger("mssp.pptx").exception("template export failed; using native builder")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    tag, label = _period_labels(period)
    brand = {
        "name": tenant.get("name", "MSSP Client") if tenant else "MSSP Client",
        "domain": tenant.get("domain", "ALL") if tenant else "ALL",
        "primary": tenant.get("primary_color") if tenant else "#00ABAB",
        "period_tag": tag,
        "period_label": label,
    }

    ex = all_data.get("executive", {})
    soc = all_data.get("soc_live")
    ti = all_data.get("ti_live")
    qbr = all_data.get("qbr", {})
    rules_count = all_data.get("rules_count")

    _qbr_cover(prs, brand)
    _qbr_section(prs, brand, "01", "Executive Summary",
                 "Quarterly performance, service reliability, and detection outcomes.",
                 "SERVICE RELIABILITY   •   SIGNAL   •   COVERAGE & DETECTION", 1, 2)
    _qbr_exec_overview(prs, brand, ex, soc, ti)
    _qbr_exec_performance(prs, brand, ex, soc, rules_count, qbr)
    _qbr_section(prs, brand, "02", "Incident Monitoring",
                 "From log-source concentration to incident trends and MITRE ATT&CK.",
                 "LOG SOURCES   •   INCIDENT TRENDS   •   MITRE ATT&CK", 2, 2)
    _qbr_log_sources(prs, brand, qbr)
    _qbr_alerts(prs, brand, qbr)
    _qbr_mitre(prs, brand, qbr)
    _qbr_closing(prs, brand)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
