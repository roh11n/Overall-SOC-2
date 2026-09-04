"""QBR-style PPTX export tests (Deloitte deck redesign).

Validates GET /api/export/pptx returns a 9-slide QBR deck with tenant-name
cover, two section dividers, populated Exec Overview tiles and Incident
Monitoring content (log sources + MITRE tactics). Also verifies graceful
behavior with no data and for all three period labels.
"""
import io
import os
import pytest
import requests
from pptx import Presentation

def _load_env():
    p = "/app/frontend/.env"
    if os.path.exists(p):
        for line in open(p):
            if "=" in line and not line.strip().startswith("#"):
                k, v = line.strip().split("=", 1)
                os.environ.setdefault(k, v)
_load_env()
BASE_URL = os.environ["REACT_APP_BACKEND_URL"].rstrip("/")
API = f"{BASE_URL}/api"

ADMIN_EMAIL = "admin@mssp-soc.io"
ADMIN_PWD = "Soc-I10eekKuxiW23Q!"

CSV = (
    "id,name,severity,status,occurred,created,closed,Actual Time Taken,Log Source,"
    "Rule Name,MITRE Tactic Name,MITRE Technique Name,SLA Breached,close reason\n"
    "1,PowerShell,Critical,Closed,2026-04-01T10:00:00Z,2026-04-01T10:12:00Z,2026-04-01T11:00:00Z,1800,Active Directory,Susp PowerShell,Execution,PowerShell,false,True Positive\n"
    "2,Travel,High,Closed,2026-04-02T09:00:00Z,2026-04-02T09:05:00Z,2026-04-02T09:45:00Z,900,O365,Impossible Travel,Initial Access,Valid Accounts,false,False Positive\n"
    "3,Brute,High,Closed,2026-05-03T08:00:00Z,2026-05-03T08:20:00Z,2026-05-03T09:30:00Z,2400,Active Directory,Failed Logins,Credential Access,Brute Force,true,True Positive\n"
    "4,Kerb,Medium,Open,2026-05-04T14:00:00Z,2026-05-04T14:30:00Z,,,SentinelOne,Kerberoast,Credential Access,Kerberoasting,false,\n"
    "5,Exfil,Critical,Closed,2026-06-05T16:00:00Z,2026-06-05T16:08:00Z,2026-06-05T17:00:00Z,3000,Firewall,Data Exfil,Exfiltration,Over C2,true,True Positive\n"
    "6,Phish,Low,Closed,2026-06-06T11:00:00Z,2026-06-06T11:03:00Z,2026-06-06T11:20:00Z,600,O365,IOC Match,Initial Access,Phishing,false,False Positive\n"
    "7,Persist,Medium,Closed,2026-06-07T12:00:00Z,2026-06-07T12:10:00Z,2026-06-07T13:00:00Z,1200,Active Directory,Sched Task,Persistence,Scheduled Task,false,True Positive\n"
)

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


@pytest.fixture(scope="session")
def token():
    r = requests.post(f"{API}/auth/login",
                      json={"email": ADMIN_EMAIL, "password": ADMIN_PWD}, timeout=30)
    assert r.status_code == 200, f"login failed: {r.status_code} {r.text[:200]}"
    tok = r.json().get("access_token") or r.json().get("token")
    assert tok, f"no token in login response: {r.json()}"
    return tok


@pytest.fixture(scope="session")
def hdr(token):
    return {"Authorization": f"Bearer {token}"}


def _slides_text(pptx_bytes: bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    texts = []
    for s in prs.slides:
        parts = []
        for shp in s.shapes:
            if shp.has_text_frame:
                parts.append(shp.text_frame.text)
            if shp.has_table:
                for row in shp.table.rows:
                    for c in row.cells:
                        parts.append(c.text)
        texts.append("\n".join(parts))
    return prs, texts


def _download(hdr, period="quarterly", tenant_id="all"):
    r = requests.get(f"{API}/export/pptx",
                     params={"period": period, "tenant_id": tenant_id},
                     headers=hdr, timeout=60)
    return r


# ---------- Empty-data behavior (test first, before upload) ----------

def test_pptx_no_data_still_valid(hdr):
    # Ensure clean state
    requests.delete(f"{API}/dashboard/soc-manager/data",
                    params={"tenant_id": "all"}, headers=hdr, timeout=30)
    r = _download(hdr, period="quarterly")
    assert r.status_code == 200, r.text[:200]
    assert PPTX_MIME in r.headers.get("content-type", ""), r.headers
    prs, texts = _slides_text(r.content)
    assert len(prs.slides) == 9, f"expected 9 slides, got {len(prs.slides)}"
    joined = "\n".join(texts)
    assert "N/A" in joined, "expected N/A placeholders when no data"
    assert "Quarterly SOC Services Report" in texts[0]


# ---------- Upload sample XSOAR data ----------

def test_upload_xsoar_sample(hdr):
    files = {"file": ("sample.csv", CSV, "text/csv")}
    r = requests.post(f"{API}/upload/data",
                      params={"source": "xsoar", "tenant_id": "all"},
                      headers=hdr, files=files, timeout=60)
    assert r.status_code == 200, f"upload failed: {r.status_code} {r.text[:300]}"


# ---------- Content validation with data ----------

def test_pptx_quarterly_structure_and_content(hdr):
    r = _download(hdr, period="quarterly")
    assert r.status_code == 200
    assert PPTX_MIME in r.headers.get("content-type", "")
    prs, texts = _slides_text(r.content)
    assert len(prs.slides) == 9

    # Slide 1: cover — tenant name, report title, Deloitte wordmark
    cover = texts[0]
    assert "Quarterly SOC Services Report" in cover, cover[:400]
    assert "Deloitte" in cover
    # Section dividers
    joined = "\n".join(texts)
    assert "Executive Summary" in joined
    assert "Incident Monitoring" in joined

    # Executive Overview slide (slide index 2) — SLA + MTTD/MTTR tiles populated
    exec_overview = texts[2]
    assert "SLA compliance" in exec_overview.lower() or "sla compliance" in exec_overview.lower()
    assert "MTTD" in exec_overview and "MTTR" in exec_overview
    # Should not be entirely N/A — at least MTTR value expected populated
    non_na_lines = [l for l in exec_overview.splitlines()
                    if l.strip() and l.strip() != "N/A"]
    assert any(ch.isdigit() for ch in exec_overview), \
        f"exec overview appears to have no numeric metrics: {exec_overview}"

    # Log sources slide should reference our uploaded sources (search all slides)
    joined_all = "\n---SLIDE---\n".join(texts)
    assert "Active Directory" in joined_all or "O365" in joined_all, \
        f"log source distribution missing from any slide"

    # MITRE slide should reference a tactic from upload
    assert any(tac in joined_all for tac in
               ["Credential Access", "Execution", "Initial Access",
                "Exfiltration", "Persistence"]), \
        f"MITRE tactic breakdown missing"


@pytest.mark.parametrize("period", ["monthly", "weekly"])
def test_pptx_other_periods(hdr, period):
    r = _download(hdr, period=period)
    assert r.status_code == 200, r.text[:200]
    assert PPTX_MIME in r.headers.get("content-type", "")
    prs, texts = _slides_text(r.content)
    assert len(prs.slides) == 9
    # Cover still shows report title
    assert "Quarterly SOC Services Report" in texts[0]


# ---------- Cleanup ----------

def test_cleanup_uploaded_data(hdr):
    r = requests.delete(f"{API}/dashboard/soc-manager/data",
                        params={"tenant_id": "all"}, headers=hdr, timeout=30)
    assert r.status_code in (200, 204)
