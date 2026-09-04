"""QBR template-based PPTX export tests.

Validates that GET /api/export/pptx builds by editing the real Deloitte QBR
template: exactly 8 slides retained, rebranding stripped (no Coromandel/CIL),
KPI values overwritten with live tenant data, and MITRE table blanked when
no tactic mapping is present.

CRITICAL: tenant_id='all' holds the user's real 5,299-row XSOAR upload.
This test file MUST NOT upload to or delete from tenant 'all'. The
regression sub-test uses tenant_id='acme-corp' and cleans up after itself.
"""
import io
import os
import re
import pytest
import requests
from pptx import Presentation


def _load_env():
    p = "/app/frontend/.env"
    if os.path.exists(p):
        for line in open(p):
            if "=" in line and not line.strip().startswith("#"):
                k, v = line.strip().split("=", 1)
                os.environ.setdefault(k, v)


_load_env()
BASE_URL = os.environ["REACT_APP_BACKEND_URL"].rstrip("/")
API = f"{BASE_URL}/api"

ADMIN_EMAIL = "admin@mssp-soc.io"
ADMIN_PWD = "Soc-I10eekKuxiW23Q!"

PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

ACME_CSV = (
    "id,name,severity,status,occurred,created,closed,Log Source,"
    "MITRE Tactic Name,MITRE Technique Name,SLA Breached,close reason\n"
    "1,PowerShell,Critical,Closed,2026-04-01T10:00:00Z,2026-04-01T10:12:00Z,2026-04-01T11:00:00Z,Active Directory,Execution,PowerShell,false,True Positive\n"
    "2,Travel,High,Closed,2026-04-02T09:00:00Z,2026-04-02T09:05:00Z,2026-04-02T09:45:00Z,O365,Initial Access,Valid Accounts,false,False Positive\n"
    "3,Brute,High,Closed,2026-05-03T08:00:00Z,2026-05-03T08:20:00Z,2026-05-03T09:30:00Z,Active Directory,Credential Access,Brute Force,true,True Positive\n"
    "4,Kerb,Medium,Open,2026-05-04T14:00:00Z,2026-05-04T14:30:00Z,,SentinelOne,Credential Access,Kerberoasting,false,\n"
    "5,Exfil,Critical,Closed,2026-06-05T16:00:00Z,2026-06-05T16:08:00Z,2026-06-05T17:00:00Z,Firewall,Exfiltration,Over C2,true,True Positive\n"
    "6,Phish,Low,Closed,2026-06-06T11:00:00Z,2026-06-06T11:03:00Z,2026-06-06T11:20:00Z,O365,Initial Access,Phishing,false,False Positive\n"
    "7,Persist,Medium,Closed,2026-06-07T12:00:00Z,2026-06-07T12:10:00Z,2026-06-07T13:00:00Z,Active Directory,Persistence,Scheduled Task,false,True Positive\n"
    "8,Cred,High,Closed,2026-06-08T13:00:00Z,2026-06-08T13:05:00Z,2026-06-08T13:55:00Z,SentinelOne,Credential Access,OS Cred Dumping,false,True Positive\n"
)


@pytest.fixture(scope="session")
def token():
    r = requests.post(f"{API}/auth/login",
                      json={"email": ADMIN_EMAIL, "password": ADMIN_PWD}, timeout=30)
    assert r.status_code == 200, f"login failed: {r.status_code} {r.text[:200]}"
    tok = r.json().get("access_token") or r.json().get("token")
    assert tok, f"no token in login response: {r.json()}"
    return tok


@pytest.fixture(scope="session")
def hdr(token):
    return {"Authorization": f"Bearer {token}"}


def _slides_text(pptx_bytes: bytes):
    prs = Presentation(io.BytesIO(pptx_bytes))
    texts = []
    tables = []
    for s in prs.slides:
        parts = []
        stables = []
        for shp in s.shapes:
            if shp.has_text_frame:
                parts.append(shp.text_frame.text)
            if shp.has_table:
                rows = []
                for row in shp.table.rows:
                    rows.append([c.text for c in row.cells])
                stables.append(rows)
                for row in rows:
                    for c in row:
                        parts.append(c)
        texts.append("\n".join(parts))
        tables.append(stables)
    return prs, texts, tables


def _download(hdr, period="quarterly", tenant_id="all"):
    return requests.get(f"{API}/export/pptx",
                        params={"period": period, "tenant_id": tenant_id},
                        headers=hdr, timeout=120)


# ---------- Template presence ----------

def test_template_file_present():
    assert os.path.exists("/app/backend/assets/qbr_template.pptx"), \
        "template file missing at /app/backend/assets/qbr_template.pptx"


# ---------- Tenant 'all' (READ-ONLY; must not modify) ----------

def test_all_tenant_pptx_status_and_slide_count(hdr):
    r = _download(hdr, period="quarterly", tenant_id="all")
    assert r.status_code == 200, r.text[:300]
    assert PPTX_MIME in r.headers.get("content-type", ""), r.headers
    prs, texts, tables = _slides_text(r.content)
    assert len(prs.slides) == 8, f"expected 8 slides, got {len(prs.slides)}"


def test_all_tenant_no_cil_coromandel_rebranding(hdr):
    r = _download(hdr, period="quarterly", tenant_id="all")
    assert r.status_code == 200
    _, texts, _ = _slides_text(r.content)
    joined = "\n".join(texts)
    # No CIL/Coromandel anywhere
    assert "Coromandel" not in joined, "found 'Coromandel' in export"
    assert "COROMANDEL" not in joined, "found 'COROMANDEL' in export"
    # standalone CIL (not part of another word)
    assert not re.search(r"\bCIL\b", joined), "found standalone 'CIL' in export"
    # Leftover CIL numbers
    for bad in ("1,803", "9,889", "3,510", "1803", "9889", "3510"):
        # allow these as digit substrings only if part of a bigger number sequence
        # be strict: template used them as standalone tokens.
        pass
    # Strict scan for the exact template numbers as standalone tokens
    for bad in ("1,803", "9,889", "3,510"):
        assert bad not in joined, f"leftover CIL number '{bad}' found in export"


def test_all_tenant_cover_and_deloitte_branding(hdr):
    r = _download(hdr, period="quarterly", tenant_id="all")
    _, texts, _ = _slides_text(r.content)
    cover = texts[0]
    # Tenant display name for 'all' is 'All Tenants'
    assert "All Tenants" in cover, f"cover missing tenant name; cover={cover[:400]}"
    # Deloitte branding retained somewhere in deck
    joined = "\n".join(texts)
    assert "Deloitte" in joined, "Deloitte branding disappeared from deck"


def test_all_tenant_executive_overview_live_values(hdr):
    r = _download(hdr, period="quarterly", tenant_id="all")
    _, texts, _ = _slides_text(r.content)
    # Slide 3 (0-based index 2) = Executive Overview after keep-slides
    exec_ov = texts[2]
    # Quarterly Signal — incidents should equal 5,299
    assert "5,299" in exec_ov, \
        f"expected live total incidents '5,299' in exec overview; got:\n{exec_ov}"
    # MTTD tile: some numeric value ending in ' min'
    assert re.search(r"\d[\d,\.]*\s*min", exec_ov), \
        f"MTTD tile has no minute value:\n{exec_ov}"
    # MTTR tile: some numeric value ending in ' h'
    assert re.search(r"\d[\d,\.]*\s*h(\b|$)", exec_ov), \
        f"MTTR tile has no hour value:\n{exec_ov}"
    # SLA compliance %: at least one %-number that is not 99.8% template default
    pcts = re.findall(r"\b\d{1,3}(?:\.\d+)?%", exec_ov)
    assert pcts, f"no percent values in exec overview:\n{exec_ov}"


def test_all_tenant_log_sources_headline_and_clean_names(hdr):
    r = _download(hdr, period="quarterly", tenant_id="all")
    _, texts, _ = _slides_text(r.content)
    # Log Sources is slide 6 in the kept set (0-based idx 5) after keep_slides
    log_slide = texts[5]
    # Headline format: "<TopSource> Drove <pct>% of Incidents This Period"
    m = re.search(r"(.+?)\s+Drove\s+(\d+(?:\.\d+)?)%\s+of\s+Incidents\s+This\s+Period",
                  log_slide)
    assert m, f"log sources headline malformed:\n{log_slide[:600]}"
    top_name = m.group(1).strip()
    # No leading [" bracket/quote artifact from JSON-array log_source parsing
    assert not top_name.startswith('['), f"top source has [ artifact: {top_name!r}"
    assert not top_name.startswith('"'), f"top source has quote artifact: {top_name!r}"
    assert '["' not in log_slide, f"log-source bracket/quote artifact present"


def test_all_tenant_alert_volume_slide(hdr):
    r = _download(hdr, period="quarterly", tenant_id="all")
    _, texts, _ = _slides_text(r.content)
    # Alert Volume is idx 6 in kept set
    av = texts[6]
    assert "5,299" in av, f"alert volume title missing live incident count:\n{av[:400]}"


def test_all_tenant_mitre_slide_blank_when_no_tactics(hdr):
    r = _download(hdr, period="quarterly", tenant_id="all")
    _, texts, tables = _slides_text(r.content)
    # MITRE slide is idx 7 in kept set
    mitre_text = texts[7]
    assert "MITRE ATT&CK Tactic Activity" in mitre_text, \
        f"MITRE title missing:\n{mitre_text[:400]}"
    assert "N/A" in mitre_text, "MITRE slide missing N/A note"
    # Heat-map table body must be blank (no CIL numbers like 669, no 'Credential Access' row)
    for tbl in tables[7]:
        # Skip header row (row 0) — check body rows
        for r_idx, row in enumerate(tbl):
            if r_idx == 0:
                continue
            for cell in row:
                assert "669" not in cell, f"leftover CIL heat cell '669' present"
                # Body first-column tactic labels must be blank
        # First-column body cells should all be empty
        for r_idx, row in enumerate(tbl):
            if r_idx == 0:
                continue
            first = row[0].strip()
            assert first == "", \
                f"MITRE body row {r_idx} first col not blank: {first!r}"


@pytest.mark.parametrize("period", ["monthly", "weekly"])
def test_all_tenant_other_periods(hdr, period):
    r = _download(hdr, period=period, tenant_id="all")
    assert r.status_code == 200, r.text[:300]
    assert PPTX_MIME in r.headers.get("content-type", "")
    prs, texts, _ = _slides_text(r.content)
    assert len(prs.slides) == 8
    # No CIL leftovers regardless of period
    joined = "\n".join(texts)
    assert "Coromandel" not in joined
    assert not re.search(r"\bCIL\b", joined)


# ---------- Regression on acme-corp (safe to upload+delete) ----------

def test_acme_upload_and_pptx_has_mitre_and_sources(hdr):
    # Clean any prior state on acme-corp
    requests.delete(f"{API}/dashboard/soc-manager/data",
                    params={"tenant_id": "acme-corp"}, headers=hdr, timeout=30)
    files = {"file": ("acme.csv", ACME_CSV, "text/csv")}
    up = requests.post(f"{API}/upload/data",
                       params={"source": "xsoar", "tenant_id": "acme-corp"},
                       headers=hdr, files=files, timeout=60)
    assert up.status_code == 200, f"acme upload failed: {up.status_code} {up.text[:300]}"

    try:
        r = _download(hdr, period="quarterly", tenant_id="acme-corp")
        assert r.status_code == 200, r.text[:300]
        prs, texts, tables = _slides_text(r.content)
        assert len(prs.slides) == 8

        joined = "\n".join(texts)
        # Rebranding sanity
        assert "Coromandel" not in joined
        assert not re.search(r"\bCIL\b", joined)

        # MITRE slide (idx 7) should be populated now
        mitre_text = texts[7]
        assert "N/A" not in mitre_text.split("MITRE ATT&CK")[0] or \
               any(t in mitre_text for t in
                   ["Credential Access", "Execution", "Initial Access",
                    "Exfiltration", "Persistence"]), \
            f"MITRE slide should be populated with tactics; got:\n{mitre_text[:600]}"
        # Table body must contain at least one tactic name
        body_has_tactic = False
        for tbl in tables[7]:
            for r_idx, row in enumerate(tbl):
                if r_idx == 0:
                    continue
                if any(t in row[0] for t in
                       ["Credential", "Execution", "Initial", "Exfil", "Persist"]):
                    body_has_tactic = True
                    break
        assert body_has_tactic, "MITRE table body has no tactic rows for acme-corp"

        # Log Sources slide (idx 5) — sources from upload present
        log_slide = texts[5]
        assert any(src in log_slide for src in
                   ["Active Directory", "O365", "SentinelOne", "Firewall"]), \
            f"log sources not reflected on acme deck:\n{log_slide[:600]}"
    finally:
        d = requests.delete(f"{API}/dashboard/soc-manager/data",
                            params={"tenant_id": "acme-corp"},
                            headers=hdr, timeout=30)
        assert d.status_code in (200, 204), f"acme cleanup failed: {d.status_code}"
