"""QBR PPTX built by editing the client's real deck as a template.

We load the original Deloitte QBR file, keep only the Executive Summary and
Incident Monitoring slides, overwrite every KPI value in place with this
tenant's live numbers (N/A where we don't track a metric), replace the S9
charts and the S10 MITRE table with our data, and rebrand to the tenant.
This preserves the exact layout, fonts, spacing and wording of the source.
"""
from __future__ import annotations

import io
import logging
import os
from datetime import datetime, timezone

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor

logger = logging.getLogger("mssp.pptx_template")

TEMPLATE_PATH = os.path.join(os.path.dirname(__file__), "assets", "qbr_template.pptx")

# 0-based indices of slides to keep: cover, Exec divider, Exec Overview,
# Exec Performance, Incident divider, Log Sources, Alert Volume, MITRE.
KEEP = {0, 3, 4, 5, 6, 7, 8, 9}

_MONTHS = ["JANUARY", "FEBRUARY", "MARCH", "APRIL", "MAY", "JUNE",
           "JULY", "AUGUST", "SEPTEMBER", "OCTOBER", "NOVEMBER", "DECEMBER"]

NAVY = RGBColor(0x01, 0x21, 0x69)
TEAL = RGBColor(0x00, 0xAB, 0xAB)
TEAL_LT = RGBColor(0x86, 0xD8, 0xD2)
GRID_LT = RGBColor(0xEE, 0xF3, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


# ---------- low-level helpers -------------------------------------------

def _set(shape, text):
    """Overwrite a shape's text, keeping the first run's formatting."""
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = str(text)
        for r in list(p.runs)[1:]:
            r._r.getparent().remove(r._r)
    else:
        p.add_run().text = str(text)
    for extra in list(tf.paragraphs)[1:]:
        extra._p.getparent().remove(extra._p)


def _set_idx(slide, idx, text):
    shapes = list(slide.shapes)
    if 0 <= idx < len(shapes):
        _set(shapes[idx], text)


def _tokens(slide, mapping):
    """Substring replace across runs (rebranding / period tokens)."""
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                for a, b in mapping.items():
                    if a in run.text:
                        run.text = run.text.replace(a, b)


def _keep_slides(prs, keep):
    lst = prs.slides._sldIdLst
    for i, sid in enumerate(list(lst)):
        if i not in keep:
            lst.remove(sid)


def _replace_chart(shape, categories, series):
    cd = CategoryChartData()
    cd.categories = list(categories)
    for name, vals in series:
        cd.add_series(name, [float(v) for v in vals])
    shape.chart.replace_data(cd)


def _heat(v, mx):
    if not v or v <= 0:
        return GRID_LT, NAVY
    t = v / mx if mx else 0
    if t > 0.66:
        return NAVY, WHITE
    if t > 0.33:
        return TEAL, WHITE
    return TEAL_LT, NAVY


def _cell(table, r, c, text, fill=None, txt=None):
    cell = table.cell(r, c)
    tf = cell.text_frame
    p = tf.paragraphs[0]
    if p.runs:
        p.runs[0].text = str(text)
        for rn in list(p.runs)[1:]:
            rn._r.getparent().remove(rn._r)
        if txt is not None:
            p.runs[0].font.color.rgb = txt
    else:
        run = p.add_run()
        run.text = str(text)
        if txt is not None:
            run.font.color.rgb = txt
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill


# ---------- formatting --------------------------------------------------

def _num(v):
    return "N/A" if v is None else f"{v:,}"


def _pct(v):
    return "N/A" if v is None else f"{v}%"


def _period(period):
    now = datetime.now(timezone.utc)
    p = (period or "monthly").lower()
    q = (now.month - 1) // 3
    if p == "quarterly":
        a, b = _MONTHS[q * 3], _MONTHS[q * 3 + 2]
        return f"{a} — {b} {now.year}", f"Q{q + 1} {now.year}"
    if p == "weekly":
        return now.strftime("WEEK OF %d %B %Y").upper(), f"WK {now.isocalendar()[1]} {now.year}"
    return now.strftime("%B %Y").upper(), now.strftime("%b %Y").upper()


# ---------- per-slide value injection -----------------------------------

def _fold_series(abm):
    """Return months + (high=critical+high, medium, low) folded to 3 series."""
    months = abm.get("months") or []
    s = abm.get("series", {})
    high = [(s.get("Critical", [0] * len(months))[i] + s.get("High", [0] * len(months))[i])
            for i in range(len(months))]
    med = list(s.get("Medium", [0] * len(months)))
    low = list(s.get("Low", [0] * len(months)))
    return months, high, med, low


def build_from_template(tenant, period, all_data, recs):
    prs = Presentation(TEMPLATE_PATH)
    slides = list(prs.slides)

    soc = all_data.get("soc_live") or {}
    s = soc.get("summary", {})
    ex = all_data.get("executive") or {}
    qbr = all_data.get("qbr") or {}
    ti = all_data.get("ti_live") or {}
    ti_live = ti.get("data_status") == "live"
    tsum = ti.get("summary", {}) if ti_live else {}

    tenant_name = (tenant or {}).get("name") or "MSSP Client"
    label, tag = _period(period)

    sla = _pct(s.get("sla_compliance_pct", ex.get("sla_compliance")))
    fp = _pct(s.get("false_positive_rate", ex.get("false_positive_rate")))
    total = s.get("total_incidents", ex.get("incidents"))
    total_s = _num(total)
    mttd = s.get("mttd_minutes")
    mttr = s.get("mttr_hours")
    mttd_s = "N/A" if mttd is None else f"{mttd} min"
    mttr_s = "N/A" if mttr is None else f"{mttr} h"
    uniq_src = _num(qbr.get("unique_log_sources"))
    uniq_det = _num(qbr.get("unique_detections"))
    adv = _num(tsum.get("total_advisories") if ti_live else None)
    ioc = _num(tsum.get("total_iocs") if ti_live else None)

    # ---- global rebrand + period on all kept slides ----
    tok = {
        "COROMANDEL INTERNATIONAL LIMITED": tenant_name.upper(),
        "Coromandel International Limited": tenant_name,
        "Coromandel": tenant_name,
        "CIL's": tenant_name + "'s", "CIL’s": tenant_name + "'s", "CIL": tenant_name,
        "/ 05": "/ 02",
        "APRIL — JUNE 2026": label, "APRIL - JUNE 2026": label,
        "Q2 2026": tag,
    }
    for i in KEEP:
        try:
            _tokens(slides[i], tok)
        except Exception:
            logger.exception("token rebrand failed on slide %s", i)

    # ---- Slide 5 (idx4): Executive Overview ----
    try:
        sl = slides[4]
        for idx, val in {
            4: sla, 7: "N/A", 10: mttd_s, 13: mttr_s,
            18: "N/A", 21: adv, 24: ioc, 27: uniq_src, 30: uniq_det,
            34: "N/A", 36: "N/A", 38: "N/A", 40: "N/A",
            44: "N/A", 46: total_s, 52: fp,
        }.items():
            _set_idx(sl, idx, val)
    except Exception:
        logger.exception("slide5 inject failed")

    # ---- Slide 6 (idx5): Executive Performance View ----
    try:
        sl = slides[5]
        headline = (
            f"Broader visibility is translating into faster, more resilient defense — "
            f"SLA {sla}, MTTR {mttr_s}."
        )
        _set_idx(sl, 4, headline)
        for idx, val in {
            22: total_s, 24: uniq_src, 26: uniq_det, 28: "N/A",
            32: "N/A", 34: "N/A", 35: "N/A",
            39: sla, 42: "N/A",
        }.items():
            _set_idx(sl, idx, val)
    except Exception:
        logger.exception("slide6 inject failed")

    # ---- Slide 8 (idx7): Log Sources ----
    try:
        sl = slides[7]
        srcs = qbr.get("log_sources") or []
        # template (name_idx, pct_idx) pairs ordered by their original % desc
        pairs = [(19, 18), (4, 3), (7, 6), (2, 21), (9, 17),
                 (8, 5), (27, 28), (24, 26), (20, 1), (23, 25)]
        if srcs:
            top = srcs[0]
            _set_idx(sl, 0, f"{top['name']} Drove {top['pct']}% of Incidents This Period")
        _set_idx(sl, 29,
                 "Incident volume by originating log source for the reporting period, "
                 "ranked from highest to lowest share.")
        for i, (nidx, pidx) in enumerate(pairs):
            if i < len(srcs):
                sc = srcs[i]
                _set_idx(sl, nidx, sc["name"])
                _set_idx(sl, pidx, f"{sc['pct']}% ({sc['count']}/{total or 0})")
            else:
                _set_idx(sl, nidx, "—")
                _set_idx(sl, pidx, "0%")
        for idx in (30, 32, 34):
            _set_idx(sl, idx, "")
    except Exception:
        logger.exception("slide8 inject failed")

    # ---- Slide 9 (idx8): Alert Volume by Month (2 native charts) ----
    try:
        sl = slides[8]
        months, high, med, low = _fold_series(qbr.get("alerts_by_month") or {})
        _set_idx(sl, 0, f"Alert Volume by Month — {total_s} Incidents This Period")
        charts = [sh for sh in sl.shapes if sh.has_chart]
        if months:
            if len(charts) >= 1:
                _replace_chart(charts[0], months,
                               [("High", high), ("Medium", med), ("Low", low)])
            if len(charts) >= 2:
                _replace_chart(charts[1], ["High", "Medium", "Low"],
                               [(m, [high[i], med[i], low[i]]) for i, m in enumerate(months)])
        else:
            _set_idx(sl, 0, "Alert Volume by Month — N/A")
            for ch in charts:
                _replace_chart(ch, ["No data"], [("High", [0]), ("Medium", [0]), ("Low", [0])])
    except Exception:
        logger.exception("slide9 inject failed")

    # ---- Slide 10 (idx9): MITRE ATT&CK heat-map table ----
    try:
        sl = slides[9]
        tbl_shape = next((sh for sh in sl.shapes if sh.has_table), None)
        tt = qbr.get("tactic_table") or []
        tmonths = qbr.get("table_months") or []
        if tbl_shape is not None:
            table = tbl_shape.table
            ncols = len(table.columns)      # 8: tactic + 6 months + total
            nrows = len(table.rows)          # 15: header + 14
            mcol = ncols - 2                 # month columns available (6)
            if tt:
                _set_idx(sl, 1, f"{tt[0]['tactic']} Led the Quarter")
                for c in range(1, 1 + mcol):
                    mi = c - 1
                    _cell(table, 0, c, tmonths[mi][:3].upper() if mi < len(tmonths) else "")
                mx = max((v for row in tt for v in row["months"].values()), default=1) or 1
                for ridx in range(1, nrows):
                    di = ridx - 1
                    if di < len(tt):
                        row = tt[di]
                        _cell(table, ridx, 0, row["tactic"][:22])
                        for c in range(1, 1 + mcol):
                            mi = c - 1
                            if mi < len(tmonths):
                                val = row["months"].get(tmonths[mi], 0)
                                fill, txt = _heat(val, mx)
                                _cell(table, ridx, c, val, fill, txt)
                            else:
                                _cell(table, ridx, c, "", WHITE, NAVY)
                        _cell(table, ridx, ncols - 1, row["total"], TEAL, WHITE)
                    else:
                        for c in range(ncols):
                            _cell(table, ridx, c, "", WHITE, NAVY)
                callouts = [(12, 13), (15, 16), (18, 19)]
                for i, (vidx, didx) in enumerate(callouts):
                    if i < len(tt):
                        _set_idx(sl, vidx, f"{tt[i]['total']:,}")
                        _set_idx(sl, didx, f"{tt[i]['tactic']} alerts recorded this period.")
                    else:
                        _set_idx(sl, vidx, "N/A")
                        _set_idx(sl, didx, "")
            else:
                # No MITRE tactic data in this upload — blank the CIL heat-map
                # rather than leave misleading numbers.
                _set_idx(sl, 1, "MITRE ATT&CK Tactic Activity")
                _set_idx(sl, 2, "N/A — the uploaded incidents contained no MITRE ATT&CK "
                                "tactic mapping, so no heat-map could be generated.")
                for ridx in range(1, nrows):
                    _cell(table, ridx, 0, "")
                    for c in range(1, ncols):
                        _cell(table, ridx, c, "", WHITE, NAVY)
                for vidx, didx in [(12, 13), (15, 16), (18, 19)]:
                    _set_idx(sl, vidx, "N/A")
                    _set_idx(sl, didx, "")
    except Exception:
        logger.exception("slide10 inject failed")

    # ---- keep only the two sections ----
    _keep_slides(prs, KEEP)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


def template_available():
    return os.path.exists(TEMPLATE_PATH)
